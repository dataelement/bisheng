#!/usr/bin/env python
"""Read a .pptx back and report what is actually in it.

Replaces two tools the official Anthropic pptx skill relies on but that do not
exist in the BiSheng code interpreter: ``markitdown`` (content dump) and
``office/validate.py`` (which needs ``defusedxml``).

Two sections are printed:

1. 内容 — every slide's text plus speaker notes, so content QA (typos, wrong
   order, missing sections) can be done without opening the file.
2. 体检 — geometric and editorial checks that approximate a visual review:
   text overflow, off-canvas shapes, overlap, tiny fonts, thin margins,
   leftover placeholder text, text-only slides.

Overflow is an *estimate*: PowerPoint's real line breaking depends on the font
metrics of the machine that opens the deck, which are unknowable here. The
estimate is deliberately tuned to under-report mild cases and only shout when a
box is clearly too small.

Usage:
    python skills/bisheng-pptx/scripts/inspect_deck.py output/deck.pptx
    python skills/bisheng-pptx/scripts/inspect_deck.py output/deck.pptx --checks-only

Always exits 0, including on a fatal error — see the comment in main().
"""

from __future__ import annotations

import argparse
import math
import os
import re
import sys
import traceback
from dataclasses import dataclass

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

# --- tunables -------------------------------------------------------------
#
# House rule for every threshold below: the checker sits ONE NOTCH LOOSER than
# the authoring spec (SKILL.md §7 and references/design-zh.md), and only shouts
# on a clear violation. Rationale: an ERROR here is a hard delivery gate, so a
# threshold set at the spec value would block decks that are merely at the edge
# of the recommended range (e.g. design-zh allows 10-12pt captions — an ERROR at
# "<12pt" would make a compliant deck unshippable, and the model would learn to
# ignore the report). The spec stays the authority on how to lay a slide out;
# this file only catches what is visibly broken.
#
# Every constant names the clause it derives from. Drift between the two used to
# be invisible precisely because that anchor was missing — keep it, and keep the
# Finding message consistent with the number that actually fires.

DEFAULT_FONT_PT = 18.0  # technical fallback when the size is inherited and unreadable here
DEFAULT_LINE_SPACING = 1.2  # technical fallback; design-zh section 4 recommends 1.2-1.4

# ANCHOR: SKILL.md section 7, "text must never overflow its container"; design-zh
# section 6, same clause. The spec is zero tolerance; this estimate carries about
# +-10% of slack (the font metrics of the opening machine are unknowable), so WARN
# starts at +5% over the box and ERROR at +30%.
OVERFLOW_WARN_RATIO = 1.05
OVERFLOW_ERROR_RATIO = 1.30

# ANCHOR: design-zh section 4, the font-size table — body 14-18pt, chart labels and
# captions 10-12pt, "below 12pt counts as unreadable". WARN fires below the smallest
# size the spec allows for ANY role (10pt), so a legitimate 10-12pt caption stays
# silent; ERROR fires only where no role is readable at all (<8pt).
MIN_FONT_WARN_PT = 10.0
MIN_FONT_ERROR_PT = 8.0

# ANCHOR: design-zh section 4, "keep >= 0.5 inch of margin on all four sides" = 36pt.
# Fires at 0.3in, i.e. only when the shape visibly hugs the edge rather than merely
# sitting below the recommendation.
MARGIN_WARN_PT = 21.6  # 0.3 inch

# ANCHOR: design-zh section 6 lists overlapping/clipped text among the "obviously AI"
# defects, but gives no number. Text over text is a defect at a much lower threshold
# than text over a shape: a big stat number sitting 30% into its own caption already
# renders unreadable.
OVERLAP_WARN_RATIO = 0.15
OVERLAP_ERROR_RATIO = 0.30

# ANCHOR: design-zh section 6.0, "not a single decorative rule", repeated in SKILL.md
# section 7. A "decoration" is a thin rule (timeline axis, divider) or a small mark
# (node dot, icon). Text running across one of those is a real, visible defect —
# unlike text sitting on a card background, which is the intended design.
DECOR_THIN_PT = 20.0
DECOR_SMALL_AREA_PT = 2600.0  # ≈ 0.5 in²
DECOR_HIT_RATIO = 0.10

PLACEHOLDER_PATTERNS = [
    (re.compile(r"\bx{3,}\b", re.I), "XXX 占位"),
    (re.compile(r"lorem|ipsum", re.I), "lorem ipsum"),
    (re.compile(r"\bTODO\b|\bTBD\b", re.I), "TODO/TBD"),
    (re.compile(r"待填|待补充|此处填写|请输入|占位符"), "中文占位提示"),
    (re.compile(r"\{\{.*?\}\}|\[insert|【\s*】"), "模板变量未替换"),
    (re.compile(r"单击此处|点击此处添加|Click to edit", re.I), "PowerPoint 默认提示语"),
]

EMU_PER_PT = 12700


@dataclass
class Finding:
    level: str  # ERROR / WARN / INFO
    slide: int
    where: str
    message: str

    def render(self) -> str:
        loc = f"第{self.slide}页" if self.slide else "整体"
        where = f" 「{self.where}」" if self.where else ""
        return f"{self.level:<5} {loc}{where} {self.message}"


# --- geometry helpers -----------------------------------------------------


def _pt(value) -> float:
    """EMU (int or pptx Length) -> points."""
    if value is None:
        return 0.0
    return float(value) / EMU_PER_PT


def _char_em(ch: str) -> float:
    """Rough advance width of one character, in em units of the font size."""
    o = ord(ch)
    if (
        0x4E00 <= o <= 0x9FFF  # CJK unified
        or 0x3400 <= o <= 0x4DBF  # extension A
        or 0xF900 <= o <= 0xFAFF  # compatibility
        or 0x3000 <= o <= 0x303F  # CJK punctuation
        or 0xFF01 <= o <= 0xFF60  # fullwidth forms
        or 0xFFE0 <= o <= 0xFFE6
        or 0x3040 <= o <= 0x30FF  # kana
    ):
        return 1.0
    if ch == " ":
        return 0.28
    if ch in "iljItfr().,;:'\"|!":
        return 0.32
    if ch.isupper() or ch.isdigit():
        return 0.60
    return 0.53


def _text_em_width(text: str) -> float:
    return sum(_char_em(c) for c in text)


def _resolve_size_pt(run, para, shape, layout_sizes: dict) -> float:
    """Effective font size for a run, walking the inheritance chain we can see."""
    if run is not None and run.font.size is not None:
        return run.font.size.pt
    if para.font.size is not None:
        return para.font.size.pt
    if shape.is_placeholder:
        idx = shape.placeholder_format.idx
        if idx in layout_sizes:
            return layout_sizes[idx]
    return DEFAULT_FONT_PT


def _layout_placeholder_sizes(slide) -> dict:
    """idx -> font size (pt) declared on the layout, when it declares one."""
    sizes = {}
    try:
        placeholders = slide.slide_layout.placeholders
    except (AttributeError, KeyError):
        return sizes
    for ph in placeholders:
        if not ph.has_text_frame:
            continue
        for para in ph.text_frame.paragraphs:
            size = para.font.size
            if size is None:
                for run in para.runs:
                    if run.font.size is not None:
                        size = run.font.size
                        break
            if size is not None:
                sizes[ph.placeholder_format.idx] = size.pt
                break
    return sizes


def _para_line_height_pt(para, size_pt: float) -> float:
    spacing = para.line_spacing
    if spacing is None:
        return size_pt * DEFAULT_LINE_SPACING
    if isinstance(spacing, (int, float)):
        return size_pt * float(spacing)
    return _pt(spacing)  # Length: exact leading


def _estimate_text_height_pt(shape, layout_sizes: dict) -> tuple[float, float]:
    """(estimated content height, widest single line width) in points."""
    tf = shape.text_frame
    width_pt = _pt(shape.width)
    avail_w = width_pt - _pt(tf.margin_left) - _pt(tf.margin_right)
    wrap = tf.word_wrap is not False

    total = _pt(tf.margin_top) + _pt(tf.margin_bottom)
    widest = 0.0
    for para in tf.paragraphs:
        runs = list(para.runs)
        text = "".join(r.text for r in runs)
        sizes = [_resolve_size_pt(r, para, shape, layout_sizes) for r in runs]
        size_pt = max(sizes) if sizes else _resolve_size_pt(None, para, shape, layout_sizes)

        em_width = _text_em_width(text)
        line_w = em_width * size_pt
        widest = max(widest, line_w)

        if not text:
            lines = 1
        elif not wrap or avail_w <= 0:
            lines = 1
        else:
            lines = max(1, math.ceil(line_w / avail_w))

        total += lines * _para_line_height_pt(para, size_pt)
        total += _pt(para.space_before) + _pt(para.space_after)
    return total, widest


def _effective_text_box(shape, content_h: float, auto_grow: bool) -> tuple[float, float, float, float]:
    """The band the text actually occupies, for overlap testing.

    A title placeholder is often far taller than its one line of text; using the
    raw frame would make every deck look like it has overlapping shapes. An
    auto-grow box is the opposite case — it occupies more than its stored height.
    """
    left, top = _pt(shape.left), _pt(shape.top)
    width, height = _pt(shape.width), _pt(shape.height)
    if auto_grow and content_h > height:
        return left, top, width, content_h
    if content_h >= height:
        return left, top, width, height
    anchor = shape.text_frame.vertical_anchor
    if anchor == MSO_ANCHOR.MIDDLE:
        top += (height - content_h) / 2
    elif anchor == MSO_ANCHOR.BOTTOM:
        top += height - content_h
    return left, top, width, content_h


def _rect_overlap(a, b) -> float:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    dx = min(ax + aw, bx + bw) - max(ax, bx)
    dy = min(ay + ah, by + bh) - max(ay, by)
    if dx <= 0 or dy <= 0:
        return 0.0
    return dx * dy


# --- traversal ------------------------------------------------------------


def _iter_shapes(shapes, prefix: str = ""):
    """Flatten groups; yields (shape, display_name)."""
    for shape in shapes:
        name = f"{prefix}{shape.name}"
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes, prefix=f"{name}/")
        else:
            yield shape, name


def _shape_text(shape) -> str:
    if shape.has_text_frame:
        return shape.text_frame.text
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            rows.append(" | ".join(cell.text.replace("\n", " ") for cell in row.cells))
        return "\n".join(rows)
    return ""


_MEDIA_TYPES = (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.TABLE)


def _is_visual(shape) -> bool:
    if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.TABLE):
        return True
    if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return True
    return False


# --- checks ---------------------------------------------------------------


def check_slide(slide, index: int, slide_w: float, slide_h: float) -> list[Finding]:
    findings: list[Finding] = []
    layout_sizes = _layout_placeholder_sizes(slide)
    text_boxes = []  # (rect, name) for overlap
    decorations = []  # (rect, name) thin rules / small marks text must not cross
    has_visual = False

    for shape, name in _iter_shapes(slide.shapes):
        left, top = _pt(shape.left), _pt(shape.top)
        width, height = _pt(shape.width), _pt(shape.height)
        full_bleed = left <= 1 and top <= 1 and left + width >= slide_w - 1 and top + height >= slide_h - 1

        if _is_visual(shape):
            has_visual = True

        has_text = shape.has_text_frame and shape.text_frame.text.strip()
        content_h = 0.0
        widest = 0.0
        auto_grow = False
        if has_text:
            content_h, widest = _estimate_text_height_pt(shape, layout_sizes)
            # add_textbox() writes <a:spAutoFit/> by default, so most generated
            # boxes are auto-grow: the text is never clipped, but the box swells
            # downward at render time and can push past the canvas or onto the
            # element below. Judge those by the grown height, not the stored one.
            auto_grow = shape.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        effective_h = max(height, content_h) if auto_grow else height

        # off-canvas (using the height the shape will actually occupy)
        if not full_bleed:
            over_r = (left + width) - slide_w
            over_b = (top + effective_h) - slide_h
            worst = max(-left, -top, over_r, over_b)
            if worst > 1:
                grown = " 且文字会把框撑高" if auto_grow and effective_h > height + 1 else ""
                findings.append(
                    Finding(
                        "ERROR",
                        index,
                        name,
                        f"超出画布 {worst:.0f}pt{grown}（画布 {slide_w:.0f}×{slide_h:.0f}pt）→ 移回画布内",
                    )
                )

        if not has_text and not full_bleed and shape.shape_type not in _MEDIA_TYPES:
            thin = min(width, height) < DECOR_THIN_PT
            small = width * height < DECOR_SMALL_AREA_PT
            if (thin or small) and width > 0 and height > 0:
                decorations.append(((left, top, width, height), name))

        if not shape.has_text_frame:
            continue

        if not has_text:
            if shape.is_placeholder:
                findings.append(Finding("WARN", index, name, "占位符是空的 → 填内容或删掉这个占位符"))
            continue

        text = shape.text_frame.text

        if height > 0 and not auto_grow:
            ratio = content_h / height
            if ratio >= OVERFLOW_ERROR_RATIO:
                findings.append(
                    Finding(
                        "ERROR",
                        index,
                        name,
                        f"文字预计需要 {content_h:.0f}pt，框高只有 {height:.0f}pt"
                        f"（超出 {(ratio - 1) * 100:.0f}%）→ 精简文字 / 降字号 / 加高框",
                    )
                )
            elif ratio >= OVERFLOW_WARN_RATIO:
                findings.append(
                    Finding(
                        "WARN",
                        index,
                        name,
                        f"文字可能放不下（估算 {content_h:.0f}pt vs 框高 {height:.0f}pt）→ 复核这一处",
                    )
                )
        elif height > 0 and auto_grow and content_h > height * OVERFLOW_WARN_RATIO:
            findings.append(
                Finding(
                    "WARN",
                    index,
                    name,
                    f"框设了自动调整高度，文字会把它从 {height:.0f}pt 撑到约 {content_h:.0f}pt"
                    f"（向下多占 {content_h - height:.0f}pt）→ 确认没有压到下方元素",
                )
            )

        if shape.text_frame.word_wrap is False and widest > width:
            findings.append(
                Finding(
                    "ERROR",
                    index,
                    name,
                    f"关闭了自动换行且单行宽 {widest:.0f}pt > 框宽 {width:.0f}pt → 文字会被截断",
                )
            )

        # font size
        smallest = None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                size = _resolve_size_pt(run, para, shape, layout_sizes)
                smallest = size if smallest is None else min(smallest, size)
        if smallest is not None:
            if smallest < MIN_FONT_ERROR_PT:
                findings.append(
                    Finding(
                        "ERROR",
                        index,
                        name,
                        f"字号 {smallest:.1f}pt 低于可读下限 {MIN_FONT_ERROR_PT:.0f}pt，投影上完全看不清"
                        f" → 正文改到 14-18pt，注释类不低于 10pt",
                    )
                )
            elif smallest < MIN_FONT_WARN_PT:
                findings.append(
                    Finding(
                        "WARN",
                        index,
                        name,
                        f"字号 {smallest:.1f}pt 低于注释类下限 {MIN_FONT_WARN_PT:.0f}pt"
                        f" → 正文改到 14-18pt，注释类 10-12pt",
                    )
                )

        # margins from the canvas edge
        if not full_bleed:
            gaps = [left, top, slide_w - (left + width), slide_h - (top + height)]
            tight = min(gaps)
            if 0 <= tight < MARGIN_WARN_PT:
                findings.append(
                    Finding(
                        "WARN",
                        index,
                        name,
                        f"离画布边缘只有 {tight:.0f}pt，已低于贴边报警线 {MARGIN_WARN_PT:.0f}pt(0.3in)"
                        f" → 往里挪，规范是四周留 ≥36pt(0.5in)",
                    )
                )

        # leftover placeholder copy
        for pattern, label in PLACEHOLDER_PATTERNS:
            if pattern.search(text):
                findings.append(Finding("ERROR", index, name, f"残留{label}：{text.strip()[:40]!r} → 替换成真实内容"))
                break

        if not full_bleed:
            text_boxes.append((_effective_text_box(shape, content_h, auto_grow), name))

    # overlap between text bands
    for i in range(len(text_boxes)):
        for j in range(i + 1, len(text_boxes)):
            rect_a, name_a = text_boxes[i]
            rect_b, name_b = text_boxes[j]
            inter = _rect_overlap(rect_a, rect_b)
            if inter <= 0:
                continue
            smaller = min(rect_a[2] * rect_a[3], rect_b[2] * rect_b[3])
            if smaller <= 0:
                continue
            ratio = inter / smaller
            if ratio >= OVERLAP_ERROR_RATIO:
                findings.append(
                    Finding("ERROR", index, f"{name_a} × {name_b}", f"文字区域重叠 {ratio * 100:.0f}% → 会互相压字")
                )
            elif ratio >= OVERLAP_WARN_RATIO:
                findings.append(
                    Finding("WARN", index, f"{name_a} × {name_b}", f"文字区域重叠 {ratio * 100:.0f}% → 复核间距")
                )

    # text running across a rule or a node dot
    for text_rect, text_name in text_boxes:
        for decor_rect, decor_name in decorations:
            inter = _rect_overlap(text_rect, decor_rect)
            decor_area = decor_rect[2] * decor_rect[3]
            if decor_area <= 0 or inter / decor_area < DECOR_HIT_RATIO:
                continue
            # Text sitting INSIDE the mark (a number in a circle, a label on a
            # badge) is the intended design, not a collision. Only text that
            # merely crosses the shape is a defect.
            text_area = text_rect[2] * text_rect[3]
            if text_area > 0 and inter / text_area > 0.7:
                continue
            findings.append(
                Finding(
                    "WARN",
                    index,
                    f"{text_name} × {decor_name}",
                    "文字压在装饰线/小图形上 → 把文字移开，或让线/点避开这段文字",
                )
            )

    # The cover is legitimately text-only; every other slide should carry something visual.
    if not has_visual and index > 1:
        findings.append(Finding("INFO", index, "", "整页只有文字，没有图片/图表/表格/图形 → 考虑加一个视觉元素"))

    return findings


# --- output ---------------------------------------------------------------


def dump_content(prs) -> None:
    print("=== 内容 ===")
    for i, slide in enumerate(prs.slides, start=1):
        try:
            layout_name = slide.slide_layout.name
        except (AttributeError, KeyError):
            layout_name = "?"
        print(f"\n[第{i}页] 版式: {layout_name}")
        for shape, _name in _iter_shapes(slide.shapes):
            text = _shape_text(shape)
            if not text.strip():
                continue
            for line in text.splitlines():
                if line.strip():
                    print(f"  · {line.strip()}")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                print(f"  [备注] {notes}")


def main() -> int:
    parser = argparse.ArgumentParser(description="Dump and health-check a .pptx")
    parser.add_argument("pptx", help="路径，相对工作区根，例如 output/deck.pptx")
    parser.add_argument("--checks-only", action="store_true", help="只输出体检，不 dump 内容")
    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        print(f"[FATAL] 文件不存在: {args.pptx}")
        print("请用相对工作区根的路径，例如 output/deck.pptx（不要带前导斜杠，也不要用宿主机绝对路径）。")
        return 0

    # Always exit 0, including on a crash: on a non-zero exit the code
    # interpreter returns stderr and throws stdout away, which would delete this
    # whole report. A traceback printed to stdout is far more useful than that.
    try:
        run_report(args.pptx, checks_only=args.checks_only)
    except Exception:
        print("[FATAL] 读取 / 体检 .pptx 时出错：")
        traceback.print_exc(file=sys.stdout)
        print("常见原因：文件损坏或写了一半、是加密文档、其实是老式 .ppt（python-pptx 只认 .pptx）。")
        print("→ 重新跑一次构建脚本生成完整的 .pptx，再跑本脚本；老式 .ppt 请让用户另存为 .pptx。")
    return 0


def run_report(path: str, checks_only: bool) -> None:
    prs = Presentation(path)
    slide_w, slide_h = _pt(prs.slide_width), _pt(prs.slide_height)

    if not checks_only:
        dump_content(prs)

    findings: list[Finding] = []
    slide_count = 0
    for i, slide in enumerate(prs.slides, start=1):
        slide_count = i
        findings.extend(check_slide(slide, i, slide_w, slide_h))

    print(f"\n=== 体检 ===  ({slide_count} 页, 画布 {slide_w:.0f}×{slide_h:.0f}pt)")
    order = {"ERROR": 0, "WARN": 1, "INFO": 2}
    for finding in sorted(findings, key=lambda f: (order[f.level], f.slide)):
        print(finding.render())

    errors = sum(1 for f in findings if f.level == "ERROR")
    warns = sum(1 for f in findings if f.level == "WARN")
    infos = sum(1 for f in findings if f.level == "INFO")
    print(f"\n合计: {errors} ERROR / {warns} WARN / {infos} INFO")
    if errors:
        print(f"结论: 不通过 —— 有 {errors} 项必须修复。改构建脚本重新生成，再跑一次本脚本。")
    else:
        print("结论: 通过 —— 无必须修复项。WARN 逐条复核后即可交付。")
    print("提示：体检阈值比 SKILL.md §7 / design-zh 的规范松一档，只在明显违规时出声；")
    print("      没报 ERROR 不等于完全合规，排版仍以规范为准。溢出为估算值，不必为 WARN 反复调参。")


if __name__ == "__main__":
    sys.exit(main())
