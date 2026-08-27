"""Regression guard for the three built-in skill bundles' checker scripts.

These scripts are the only mechanical quality gate the Linsight model has when it
builds a .pptx / .xlsx / .docx, and they were shipped with zero tests: a threshold
tweak, a regex edit, or an accidental non-zero exit could silently break them and
nothing in CI would notice.

Three properties are locked down here, each for a failure mode that has actually
happened in this bundle family:

* **Bidirectional.** A bad sample must be *rejected* (``结论: 不通过`` plus the
  specific rule's wording) AND a good sample must be *accepted* (``结论: 通过``).
  Testing only the bad direction lets a checker degrade into "always shouts",
  which is worse than no checker — the model learns to ignore the report.
* **``returncode == 0``, always**, including on a corrupt file. The BiSheng code
  executor is either/or: on a non-zero exit it returns *stderr* and throws the
  whole stdout report away. A script that "correctly" exits 1 on findings makes
  its own report invisible.
* **No dependency on LibreOffice.** The recalc/render scripts are exercised with
  ``soffice`` stripped from ``PATH``, asserting the documented Chinese fallback
  text rather than a crash. CI must never need LibreOffice installed.

Samples are built with python-pptx / openpyxl / python-docx, which are backend
dependencies, and the .docx good sample is built with the bundle's own
``docx_helpers`` — "the helper's own output passes the bundle's own checker" is
the contract that matters.
"""

from __future__ import annotations

import ast
import importlib.util
import os
import re
import subprocess
import sys
from pathlib import Path

import pytest
from docx import Document
from docx.shared import Cm as DocxCm
from docx.shared import Pt as DocxPt
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from bisheng.linsight.domain.services.builtin_skill_seeder import BUILTIN_SKILLS_DIR

PASS = "结论: 通过"
FAIL = "结论: 不通过"


# --------------------------------------------------------------------------- #
# running a packaged script the way SKILL.md tells the model to run it
# --------------------------------------------------------------------------- #
def run_script(bundle: str, script: str, *args, cwd: Path | None = None, path_env: str | None = None):
    """Run ``skills/<bundle>/scripts/<script>`` and assert the house rule: exit 0.

    ``path_env`` replaces ``PATH`` wholesale — used to simulate a host without
    LibreOffice. ``sys.executable`` is absolute, so the interpreter still starts.
    """
    target = BUILTIN_SKILLS_DIR / bundle / "scripts" / script
    assert target.is_file(), f"{target} is missing — SKILL.md tells the model to call it"

    env = dict(os.environ)
    if path_env is not None:
        env["PATH"] = path_env

    proc = subprocess.run(
        [sys.executable, str(target), *[str(a) for a in args]],
        capture_output=True,
        text=True,
        cwd=str(cwd) if cwd else None,
        env=env,
    )
    assert proc.returncode == 0, (
        f"{bundle}/{script} exited {proc.returncode}; a non-zero exit makes the code "
        f"executor discard the entire stdout report.\nstderr:\n{proc.stderr[-2000:]}"
    )
    assert proc.stdout.strip(), f"{bundle}/{script} printed nothing to stdout"
    return proc


def assert_mentions(stdout: str, *needles: str) -> None:
    missing = [n for n in needles if n not in stdout]
    assert not missing, f"report never mentioned {missing}\n--- report ---\n{stdout}"


# --------------------------------------------------------------------------- #
# sample builders
# --------------------------------------------------------------------------- #
def _textbox(slide, left, top, width, height, text, size_pt):
    """A textbox with the two defaults python-pptx gets wrong for us made explicit:
    ``add_textbox`` writes ``wrap="none"`` + ``<a:spAutoFit/>``, both of which the
    checker (correctly) treats as suspicious."""
    box = slide.shapes.add_textbox(Pt(left), Pt(top), Pt(width), Pt(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.text = text
    frame.paragraphs[0].runs[0].font.size = Pt(size_pt)
    return box


def _new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)  # 960x540pt
    return prs, prs.slide_layouts[6]  # blank


def build_good_pptx(path: Path) -> Path:
    prs, blank = _new_deck()

    cover = prs.slides.add_slide(blank)
    _textbox(cover, 72, 200, 816, 90, "二零二五年度电力工程经营复盘", 40)
    _textbox(cover, 72, 310, 816, 40, "规划发展部 · 二零二五年十二月", 18)

    body = prs.slides.add_slide(blank)
    _textbox(body, 72, 48, 816, 50, "全年营收与毛利", 28)
    _textbox(body, 72, 150, 380, 240, "营收同比增长两成\n毛利率提升一点四个百分点\n应收账款周转天数下降九天", 18)
    # A slide with no picture/chart/table/shape earns an INFO; a big rounded
    # rectangle is a visual and is too large to count as a decorative rule.
    body.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(520), Pt(150), Pt(360), Pt(240))

    prs.save(path)
    return path


def build_bad_pptx(path: Path) -> Path:
    prs, blank = _new_deck()

    one = prs.slides.add_slide(blank)
    _textbox(one, 72, 60, 300, 30, "这一段文字远远超出了它的容器高度" * 4, 20)  # overflow -> ERROR
    _textbox(one, 72, 300, 300, 40, "本季度增长率待填", 18)  # placeholder copy -> ERROR

    two = prs.slides.add_slide(blank)
    _textbox(two, 900, 100, 300, 60, "跑到画布外面去的文字", 18)  # off-canvas -> ERROR
    _textbox(two, 10, 300, 300, 60, "贴着画布左边缘的文字", 18)  # hugs the edge -> WARN
    _textbox(two, 400, 300, 300, 60, "极小字号的注释", 6)  # 6pt -> ERROR

    prs.save(path)
    return path


def build_good_xlsx(path: Path) -> Path:
    """No formulas on purpose: openpyxl writes formulas with no cached value, and
    that is (correctly) an ERROR until LibreOffice recalculates — which CI has no
    business needing. The formula rules get their coverage from the bad sample."""
    wb = Workbook()
    ws = wb.active
    ws.title = "营收明细"
    ws.append(["产品线", "本年营收", "去年营收", "同比增幅"])
    for name, current, previous, growth in (
        ("输配电设备", 12500000, 10800000, 0.157),
        ("电力工程服务", 8600000, 7900000, 0.089),
        ("运维托管", 4300000, 3100000, 0.387),
    ):
        ws.append([name, current, previous, growth])

    for row in ws.iter_rows(min_row=2, min_col=2, max_col=3):
        for cell in row:
            cell.number_format = "#,##0"
    for row in ws.iter_rows(min_row=2, min_col=4, max_col=4):
        for cell in row:
            cell.number_format = "0.0%"
    for letter, width in (("A", 18), ("B", 16), ("C", 16), ("D", 12)):
        ws.column_dimensions[letter].width = width
    for cell in ws[1]:
        cell.font = Font(bold=True)
    ws.freeze_panes = "A2"

    wb.save(path)
    return path


def build_bad_xlsx(path: Path) -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "假设 输入"  # a space in the name -> cross-sheet refs must be quoted
    ws["A1"] = "增长率"
    ws["B1"] = "15%"  # percentage stored as text -> ERROR
    ws["A2"] = "毛利率待填"  # leftover placeholder -> WARN
    ws["B2"] = 0.32

    calc = wb.create_sheet("测算")
    calc["A1"] = "=XLOOKUP(A2,假设 输入!A:A,假设 输入!B:B)"  # banned fn + unquoted name
    calc["A2"] = '=TEXTJOIN(",",TRUE,B1:B3)'  # missing _xlfn. prefix
    calc["A3"] = "=B1*1.15"  # hardcoded coefficient -> WARN
    calc["A4"] = "#REF!"  # baked error literal -> ERROR
    # Regression: the `_xlfn.` prefix does not un-ban a banned function, and an
    # earlier lookbehind let the `.` swallow the whole match, silencing the rule.
    calc["A5"] = "=_xlfn.SORT(B1:B3)"
    # Negative control on the same rule family: a formula that *mentions* an error
    # literal as a string is legitimate and must not be reported.
    calc["A6"] = '=IF(ISNA(MATCH(A1,B1:B3,0)),"#N/A",1)'

    wb.save(path)
    return path


def build_cjk_xlsx(path: Path) -> Path:
    """Pure-CJK sheet names referenced across sheets. Legal unquoted in Excel and
    LibreOffice, so the quoting rule must stay silent here — treating "non-ASCII"
    as "must quote" would false-positive on nearly every Chinese workbook."""
    wb = Workbook()
    assume = wb.active
    assume.title = "关键假设"
    assume["A1"], assume["B1"] = "增长率", 0.15

    detail = wb.create_sheet("收入明细2024")  # CJK + digits
    detail["A1"], detail["B1"] = "月份", "收入"
    detail["A2"], detail["B2"] = 1, 100

    total = wb.create_sheet("汇总")
    total["A1"] = "预测收入"
    total["B1"] = "=收入明细2024!B2*(1+关键假设!$B$1)"
    total["A2"] = "月份合计"
    total["B2"] = "=SUM(收入明细2024!A2:A2)"

    wb.save(path)
    return path


def _docx_helpers():
    """Import the bundle's docx_helpers by path, without polluting sys.path."""
    source = BUILTIN_SKILLS_DIR / "bisheng-docx" / "scripts" / "docx_helpers.py"
    spec = importlib.util.spec_from_file_location("builtin_docx_helpers", source)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _style_profiles():
    """Import the bundle's style_profiles by path, without polluting sys.path."""
    source = BUILTIN_SKILLS_DIR / "bisheng-docx" / "scripts" / "style_profiles.py"
    spec = importlib.util.spec_from_file_location("builtin_style_profiles", source)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _style_face(doc, style_name: str) -> str | None:
    """The ``w:eastAsia`` face a style actually carries — the only field that
    decides how a Chinese glyph renders."""
    from docx.oxml.ns import qn

    rpr = doc.styles[style_name].element.find(qn("w:rPr"))
    if rpr is None:
        return None
    rfonts = rpr.find(qn("w:rFonts"))
    return rfonts.get(qn("w:eastAsia")) if rfonts is not None else None


def build_good_docx(path: Path) -> Path:
    helpers = _docx_helpers()
    doc = Document()
    section = helpers.setup_page(doc)
    helpers.apply_chinese_defaults(doc)

    helpers.add_heading_cn(doc, "电力工程年度经营报告", 1)
    helpers.add_toc(doc)
    helpers.add_heading_cn(doc, "一、总体情况", 2)
    helpers.add_body(doc, "全年营收较上年增长两成。毛利率提升一点四个百分点。应收账款周转天数下降九天。")
    helpers.add_heading_cn(doc, "二、分产品线明细", 2)

    usable = helpers.content_width_cm(section)
    helpers.add_table(
        doc,
        ["产品线", "本年营收", "去年营收", "同比"],
        [["输配电设备", "1250", "1080", "15.7%"], ["电力工程服务", "860", "790", "8.9%"]],
        widths_cm=[usable * 0.4, usable * 0.2, usable * 0.2, usable * 0.2],
    )
    helpers.add_page_number_footer(section)

    doc.save(path)
    return path


def build_bad_docx(path: Path) -> Path:
    doc = Document()
    doc.add_heading("中文标题没有设中文字体", 1)  # style chain has no w:eastAsia -> ERROR

    run = doc.add_paragraph().add_run("这段中文只设了西文字体")
    run.font.name = "微软雅黑"  # ascii/hAnsi only -> ERROR

    small = doc.add_paragraph().add_run("这是六号小字")
    small.font.size = DocxPt(6)  # -> WARN

    doc.add_paragraph("这一段里有待填的占位内容")  # -> WARN

    table = doc.add_table(rows=2, cols=3)
    for row in table.rows:
        for cell in row.cells:
            cell.width = DocxCm(8)  # 24cm total vs a 15.2cm text width -> ERROR
            cell.text = "数据"

    doc.save(path)
    return path


def build_east_asia_only_docx(path: Path) -> Path:
    """One Chinese heading and nothing else — the single highest-value check in
    the docx bundle, isolated so a regression cannot hide behind other findings."""
    doc = Document()
    doc.add_heading("中文标题", 1)
    doc.save(path)
    return path


# --------------------------------------------------------------------------- #
# fixtures
# --------------------------------------------------------------------------- #
@pytest.fixture(scope="module")
def samples(tmp_path_factory) -> dict[str, Path]:
    root = tmp_path_factory.mktemp("builtin_skill_samples")
    built = {
        "good.pptx": build_good_pptx,
        "bad.pptx": build_bad_pptx,
        "good.xlsx": build_good_xlsx,
        "bad.xlsx": build_bad_xlsx,
        "cjk.xlsx": build_cjk_xlsx,
        "good.docx": build_good_docx,
        "bad.docx": build_bad_docx,
        "eastasia.docx": build_east_asia_only_docx,
    }
    paths = {name: builder(root / name) for name, builder in built.items()}

    # Truncated files: a half-written deliverable is the realistic corruption, and
    # every reader below must survive it with a readable stdout report.
    for stem in ("pptx", "xlsx", "docx"):
        broken = root / f"corrupt.{stem}"
        broken.write_bytes((root / f"good.{stem}").read_bytes()[:3000])
        paths[f"corrupt.{stem}"] = broken

    paths["_root"] = root
    return paths


@pytest.fixture(scope="module")
def no_soffice_path(tmp_path_factory) -> str:
    """An empty PATH entry, so shutil.which() finds no LibreOffice on any host."""
    return str(tmp_path_factory.mktemp("empty_bin"))


# --------------------------------------------------------------------------- #
# pptx
# --------------------------------------------------------------------------- #
def test_pptx_bad_deck_is_rejected_with_the_reason(samples):
    out = run_script("bisheng-pptx", "inspect_deck.py", samples["bad.pptx"], "--checks-only").stdout

    assert FAIL in out
    assert_mentions(
        out,
        "ERROR",
        "文字预计需要",  # overflow
        "残留",  # leftover placeholder copy
        "超出画布",  # off-canvas
        "字号 6.0pt 低于可读下限",  # the ERROR branch, not the softer WARN one
        "WARN",
        "离画布边缘",  # hugs the canvas edge
    )


def test_pptx_good_deck_passes(samples):
    """The checker must be satisfiable — otherwise the model learns to ignore it."""
    out = run_script("bisheng-pptx", "inspect_deck.py", samples["good.pptx"], "--checks-only").stdout

    assert PASS in out, out
    assert "合计: 0 ERROR" in out


def test_pptx_content_dump_prints_slide_text(samples):
    out = run_script("bisheng-pptx", "inspect_deck.py", samples["good.pptx"]).stdout

    assert_mentions(out, "=== 内容 ===", "二零二五年度电力工程经营复盘", "=== 体检 ===")


def test_pptx_probe_template_reads_a_deck(samples):
    out = run_script("bisheng-pptx", "probe_template.py", samples["good.pptx"]).stdout

    assert_mentions(out, "=== 画布 ===", "13.33", "16:9", "版式清单")


# --------------------------------------------------------------------------- #
# xlsx
# --------------------------------------------------------------------------- #
def test_xlsx_bad_workbook_is_rejected_with_the_reason(samples):
    out = run_script("bisheng-xlsx", "inspect_workbook.py", samples["bad.xlsx"], "--check-only").stdout

    assert FAIL in out
    assert_mentions(
        out,
        "ERROR",
        "XLOOKUP",  # banned spilling function
        "_xlfn.",  # TEXTJOIN without the prefix
        "百分比存成了文本",
        "#REF!",
        "没有缓存值",  # openpyxl formulas are never cached
        "没加单引号",  # sheet name with a space
        "用了 SORT()",  # `_xlfn.` does not un-ban a banned function
        "WARN",
        "残留占位文本",
        "硬编码了系数",
    )
    # A formula that merely *contains* "#N/A" as a string argument is legitimate;
    # only a cell whose entire content IS the literal is a baked error.
    assert "测算!A6" not in out, out


def test_xlsx_good_workbook_passes(samples):
    out = run_script("bisheng-xlsx", "inspect_workbook.py", samples["good.xlsx"], "--check-only").stdout

    assert PASS in out, out
    assert "合计: 0 ERROR / 0 WARN / 0 INFO" in out


def test_xlsx_cjk_sheet_names_do_not_false_positive(samples):
    """Regression: pure-CJK sheet names are legal unquoted, and they exist in the
    workbook — neither the quoting rule nor the unknown-sheet rule may fire."""
    out = run_script("bisheng-xlsx", "inspect_workbook.py", samples["cjk.xlsx"], "--check-only").stdout

    assert "没加单引号" not in out, out
    assert "但工作簿里没有这个表名" not in out, out
    # The only finding left is the one that is genuinely true of any freshly
    # written openpyxl workbook with formulas.
    assert "合计: 1 ERROR / 0 WARN / 0 INFO" in out, out


def test_xlsx_content_dump_prints_cells(samples):
    out = run_script("bisheng-xlsx", "inspect_workbook.py", samples["good.xlsx"], "--content-only").stdout

    assert_mentions(out, "营收明细", "A1=产品线")


def test_xlsx_recalc_degrades_without_libreoffice(samples, no_soffice_path):
    out = run_script("bisheng-xlsx", "recalc_check.py", samples["cjk.xlsx"], path_env=no_soffice_path).stdout

    assert_mentions(out, "[未重算]", "没有 soffice", "降级方案")
    # Formulas without cached values are a real delivery blocker, so this must
    # read as a failure — but in recalc's own vocabulary. `结论: 通过` / `结论: 不通过`
    # belong to inspect_workbook.py alone: SKILL.md makes `结论: 通过` the loop's
    # stop condition, and a second script emitting it would let the model finish
    # one step early, before the workbook was ever health-checked.
    assert "重算未通过" in out, out
    assert PASS not in out, out


# --------------------------------------------------------------------------- #
# docx
# --------------------------------------------------------------------------- #
def test_docx_bad_document_is_rejected_with_the_reason(samples):
    out = run_script("bisheng-docx", "inspect_docx.py", samples["bad.docx"], "--check-only").stdout

    assert FAIL in out
    assert_mentions(
        out,
        "ERROR",
        "只设了西文字体、没设 w:eastAsia",
        "超过版心",  # table wider than the text column
        "WARN",
        "字号小于 9pt",
        "残留占位文本",
    )


def test_docx_missing_east_asia_font_is_an_error(samples):
    """``doc.add_heading("中文标题", 1)`` writes no run font and the stock Heading
    style carries no CJK face, so every Chinese glyph falls back to the theme
    font. Invisible in extracted text; this is the bundle's headline check."""
    out = run_script("bisheng-docx", "inspect_docx.py", samples["eastasia.docx"], "--check-only").stdout

    assert FAIL in out, out
    # The precise ERROR wording, not just the substring "w:eastAsia" — the Normal
    # style WARN also contains that, and would keep this test green on its own.
    assert_mentions(out, "[ERROR]", "继承链里没有 w:eastAsia", "Heading 1")


def test_docx_good_document_passes(samples):
    """Built with the bundle's own docx_helpers: the helpers' output must satisfy
    the bundle's own checker, or one of the two is wrong."""
    out = run_script("bisheng-docx", "inspect_docx.py", samples["good.docx"], "--check-only").stdout

    assert PASS in out, out
    assert "合计: 0 ERROR" in out


def test_docx_content_dump_prints_paragraphs_and_tables(samples):
    out = run_script("bisheng-docx", "inspect_docx.py", samples["good.docx"], "--content-only").stdout

    assert_mentions(out, "=== 内容 ===", "电力工程年度经营报告", "## 表 1")


# --------------------------------------------------------------------------- #
# docx · 版式档 (style profiles)
# --------------------------------------------------------------------------- #
def test_docx_defaults_to_the_gongwen_profile(tmp_path):
    """Calling ``apply_chinese_defaults(doc)`` with no argument must produce a
    GB/T 9704 document.

    Two sessions shipped 公文 in 微软雅黑 because the helper hardcoded a corporate
    face and the model, told by design-zh.md to avoid "exotic" fonts, went along
    with it. The default is the whole fix — parameterising the helpers while
    leaving 微软雅黑 as the default would have changed nothing.
    """
    helpers = _docx_helpers()
    doc = Document()
    helpers.apply_chinese_defaults(doc)
    section = helpers.setup_page(doc)
    path = tmp_path / "gongwen.docx"
    doc.save(path)

    saved = Document(path)
    assert _style_face(saved, "Normal") == "仿宋_GB2312"
    assert _style_face(saved, "Heading 1") == "黑体"
    assert _style_face(saved, "Heading 2") == "楷体_GB2312"
    # 三号 throughout; the standard separates levels by face, not by size
    assert saved.styles["Normal"].font.size.pt == 16
    assert saved.styles["Heading 1"].font.size.pt == 16
    assert saved.styles["Heading 2"].font.size.pt == 16
    # 版心: 37 / 35 / 28 / 26 mm
    assert round(section.top_margin.cm, 2) == 3.7
    assert round(section.left_margin.cm, 2) == 2.8


def test_docx_modern_profile_restores_the_corporate_look(tmp_path):
    """A CV or a marketing one-pager must still be able to opt out."""
    helpers = _docx_helpers()
    doc = Document()
    helpers.apply_chinese_defaults(doc, profile="modern")
    section = helpers.setup_page(doc)
    path = tmp_path / "modern.docx"
    doc.save(path)

    saved = Document(path)
    assert _style_face(saved, "Normal") == "微软雅黑"
    assert saved.styles["Normal"].font.size.pt == 11
    assert saved.styles["Heading 1"].font.size.pt == 20
    assert round(section.top_margin.cm, 2) == 2.54


def test_setup_page_before_apply_defaults_still_lands_on_profile_margins(tmp_path):
    """SKILL.md's skeleton used to call ``setup_page`` first, and a model that
    copies the old order must not end up with 公文 margins on a modern document.

    ``apply_chinese_defaults`` therefore rewrites the margins of the sections
    already laid out, rather than trusting the call order.
    """
    helpers = _docx_helpers()
    doc = Document()
    section = helpers.setup_page(doc)  # laid out under whatever profile was active
    helpers.apply_chinese_defaults(doc, profile="modern")
    path = tmp_path / "order.docx"
    doc.save(path)

    assert round(Document(path).sections[0].top_margin.cm, 2) == 2.54
    assert round(section.left_margin.cm, 2) == 2.54


def test_profile_overrides_replace_only_the_keys_given():
    """The override dict is the seam an imported third-party 公文 skill hands its
    unit template through, so it has to merge deeply — a caller naming one face
    must not lose the size, leading and indent that came with the profile."""
    profiles = _style_profiles()

    merged = profiles.resolve_profile("gongwen", {"body": {"font": "华文中宋"}})
    assert merged["body"]["font"] == "华文中宋"
    assert merged["body"]["pt"] == 16  # untouched
    assert merged["body"]["line_pt"] == 28  # untouched
    assert merged["headings"][1]["font"] == "黑体"  # untouched

    # A bare dict is read as overrides on the default (公文) profile.
    bare = profiles.resolve_profile({"headings": {1: {"font": "方正小标宋简体"}}})
    assert bare["headings"][1]["font"] == "方正小标宋简体"
    assert bare["headings"][2]["font"] == "楷体_GB2312"
    assert bare["body"]["font"] == "仿宋_GB2312"

    # Resolving must not mutate the module-level profile.
    assert profiles.GONGWEN["body"]["font"] == "仿宋_GB2312"


def test_helpers_follow_the_active_profile_without_being_told(tmp_path):
    """Passing the profile to every call is how a document ends up 90% right:
    ``apply_chinese_defaults`` sets it once and the rest follow."""
    helpers = _docx_helpers()
    doc = Document()
    helpers.apply_chinese_defaults(doc)
    helpers.add_gongwen_title(doc, "关于××××的通知")
    helpers.add_body(doc, "正文内容。")
    helpers.add_heading_cn(doc, "一、总体情况", 1)
    helpers.add_table(doc, ["列一", "列二"], [["甲", "乙"]])
    path = tmp_path / "follow.docx"
    doc.save(path)

    from docx.oxml.ns import qn

    faces = {
        run._element.find(qn("w:rPr")).find(qn("w:rFonts")).get(qn("w:eastAsia"))
        for paragraph in Document(path).paragraphs
        for run in paragraph.runs
        if run.text.strip()
    }
    assert faces == {"方正小标宋简体", "仿宋_GB2312", "黑体"}, faces


def test_gongwen_document_built_with_generic_fonts_is_warned(tmp_path):
    """Regression for the two shipped sessions: a document that is plainly an
    official one, rendered in 微软雅黑, must be called out.

    Detection keys off the 公文 *intent* (a 「关于××的通知」title, or any of the
    four mandated faces appearing somewhere) rather than off the body face —
    keying off the body face would make the check vanish exactly when the body
    face is the thing that was replaced.
    """
    helpers = _docx_helpers()
    doc = Document()
    helpers.apply_chinese_defaults(doc, profile="modern")
    helpers.add_gongwen_title(doc, "关于下发《电力行业映射表 V2.0》的通知")
    helpers.add_body(doc, "各有关部门：", indent=False)
    helpers.add_body(doc, "现将映射表予以印发，请遵照执行。")
    path = tmp_path / "gongwen_in_yahei.docx"
    doc.save(path)

    out = run_script("bisheng-docx", "inspect_docx.py", path, "--check-only").stdout

    assert_mentions(out, "按「公文」版式档校验", "看着是公文，却用到了通用字体", "微软雅黑")
    # WARN, never ERROR: a unit template legitimately overrides the standard, and
    # an unsatisfiable checker turns the model's repair loop into an infinite one.
    assert PASS in out, out
    assert "合计: 0 ERROR" in out


def test_a_corporate_document_is_not_judged_as_gongwen(tmp_path):
    """The 公文 detector must stay narrow — a weekly report in 微软雅黑 is fine."""
    helpers = _docx_helpers()
    doc = Document()
    helpers.apply_chinese_defaults(doc, profile="modern")
    helpers.add_heading_cn(doc, "本周工作进展", 1)
    helpers.add_body(doc, "本周完成三项交付，下周进入联调。")
    path = tmp_path / "weekly.docx"
    doc.save(path)

    out = run_script("bisheng-docx", "inspect_docx.py", path, "--check-only").stdout

    assert "按「非公文」版式档校验" in out, out
    assert "通用字体" not in out, out
    assert PASS in out, out


def test_docx_render_degrades_without_libreoffice(samples, no_soffice_path):
    out = run_script(
        "bisheng-docx",
        "render_docx.py",
        samples["good.docx"],
        "--outdir",
        samples["_root"] / "preview_docx",
        path_env=no_soffice_path,
    ).stdout

    assert_mentions(out, "[跳过]", "没有 LibreOffice", "inspect_docx.py")


def test_pptx_render_degrades_without_libreoffice(samples, no_soffice_path):
    out = run_script(
        "bisheng-pptx",
        "render_deck.py",
        samples["good.pptx"],
        "--outdir",
        samples["_root"] / "preview_pptx",
        path_env=no_soffice_path,
    ).stdout

    assert_mentions(out, "找不到 soffice", "inspect_deck.py")


# --------------------------------------------------------------------------- #
# broken input — the report must survive, because a crash means an empty report
# --------------------------------------------------------------------------- #
@pytest.mark.parametrize(
    ("bundle", "script", "sample"),
    [
        ("bisheng-pptx", "inspect_deck.py", "corrupt.pptx"),
        ("bisheng-pptx", "probe_template.py", "corrupt.pptx"),
        ("bisheng-xlsx", "inspect_workbook.py", "corrupt.xlsx"),
        ("bisheng-docx", "inspect_docx.py", "corrupt.docx"),
    ],
)
def test_a_corrupt_file_still_produces_a_report_on_stdout(samples, bundle, script, sample):
    proc = run_script(bundle, script, samples[sample])  # run_script asserts exit 0

    assert "[FATAL]" in proc.stdout, proc.stdout
    assert "Traceback" in proc.stdout or "不是真的" in proc.stdout, proc.stdout
    assert not proc.stderr.strip(), f"the diagnosis leaked to stderr, where the executor drops it:\n{proc.stderr}"


def test_recalc_on_a_corrupt_file_never_crashes(samples, no_soffice_path):
    """``recalc_check.py`` is deliberately absent from the parametrisation above:
    it probes for soffice *before* opening the file, so on a host without
    LibreOffice a corrupt workbook produces the degradation notice, not a zip
    error. Both branches are exit-0 with a Chinese explanation — pin the one that
    is deterministic everywhere."""
    proc = run_script("bisheng-xlsx", "recalc_check.py", samples["corrupt.xlsx"], path_env=no_soffice_path)

    assert "[未重算]" in proc.stdout, proc.stdout
    assert not proc.stderr.strip(), proc.stderr


@pytest.mark.parametrize(
    ("bundle", "script"),
    [
        ("bisheng-pptx", "inspect_deck.py"),
        ("bisheng-pptx", "probe_template.py"),
        ("bisheng-pptx", "render_deck.py"),
        ("bisheng-xlsx", "inspect_workbook.py"),
        ("bisheng-xlsx", "recalc_check.py"),
        ("bisheng-docx", "inspect_docx.py"),
        ("bisheng-docx", "render_docx.py"),
    ],
)
def test_a_missing_file_is_explained_not_crashed(bundle, script):
    proc = run_script(bundle, script, "output/does-not-exist.bin")

    assert "[FATAL] 文件不存在" in proc.stdout, proc.stdout


# --------------------------------------------------------------------------- #
# the house rule itself
# --------------------------------------------------------------------------- #
_NONZERO_EXIT_RE = re.compile(r"(?:sys\.)?exit\(\s*[1-9]|SystemExit\(\s*[1-9]")


@pytest.mark.parametrize(
    "script",
    sorted(str(p.relative_to(BUILTIN_SKILLS_DIR)) for p in BUILTIN_SKILLS_DIR.glob("*/scripts/*.py")),
)
def test_no_script_can_exit_non_zero(script):
    """Static half of the exit-0 rule: it also covers branches no sample reaches.

    A non-zero exit makes the code executor return stderr and discard stdout, so
    the more informative the failure, the more completely it disappears.
    """
    source = (BUILTIN_SKILLS_DIR / script).read_text(encoding="utf-8")

    assert not _NONZERO_EXIT_RE.search(source), f"{script} can exit non-zero; the executor would drop its report"


@pytest.mark.parametrize(
    "script",
    sorted(str(p.relative_to(BUILTIN_SKILLS_DIR)) for p in BUILTIN_SKILLS_DIR.glob("*/scripts/*.py")),
)
def test_only_the_inspect_scripts_may_print_the_stop_string(script):
    """``结论: 通过`` is the loop's stop condition, so exactly one script owns it.

    Each SKILL.md tells the model to iterate until that string appears. A second
    script printing it — ``recalc_check.py`` is the near miss, since a clean
    recalculation feels like success — lets the model stop one step early and
    ship a file the health check never saw.
    """
    source = (BUILTIN_SKILLS_DIR / script).read_text(encoding="utf-8")
    owner = Path(script).name.startswith("inspect_")
    # Only what the script *emits* matters; prose explaining the rule is fine.
    printed = [line.strip() for line in source.splitlines() if "print(" in line and PASS in line]

    assert owner or not printed, f"{script} prints {PASS!r}, which only the inspect_* script may claim: {printed}"


_CN_FACE_LITERAL = re.compile(r"微软雅黑|仿宋|小标宋|楷体|黑体|等线|宋体")
_COMPAT_CONSTANTS = {"CN_FONT", "CN_FONT_SERIF", "CN_FONT_HEADING"}


def test_docx_helpers_read_faces_from_the_profile_not_from_literals():
    """Static half of the profile fix: no helper may name a Chinese face itself.

    The original defect was exactly this — ``add_heading_cn`` wrote
    ``CN_FONT_HEADING`` into the run, at the highest precedence there is, so a
    build script could set every style correctly and still get 微软雅黑 headings.
    The module-level compat constants may stay (build scripts reference them);
    what may not come back is a *helper* reading one, or inlining a face name.
    """
    source = (BUILTIN_SKILLS_DIR / "bisheng-docx" / "scripts" / "docx_helpers.py").read_text(encoding="utf-8")
    tree = ast.parse(source)

    offenders: list[str] = []
    for node in ast.walk(tree):
        if isinstance(node, ast.Name) and isinstance(node.ctx, ast.Load) and node.id in _COMPAT_CONSTANTS:
            offenders.append(f"line {node.lineno}: reads {node.id}")

    for func in [n for n in ast.walk(tree) if isinstance(n, ast.FunctionDef)]:
        body = func.body[1:] if ast.get_docstring(func) else func.body
        for statement in body:
            for node in ast.walk(statement):
                if isinstance(node, ast.Constant) and isinstance(node.value, str):
                    if _CN_FACE_LITERAL.search(node.value):
                        offenders.append(f"line {node.lineno}: {func.name}() inlines {node.value!r}")

    assert not offenders, "faces must come from style_profiles, not from docx_helpers: " + "; ".join(offenders)


def test_warnings_alone_do_not_block_xlsx_delivery(tmp_path):
    """WARN must never withhold the stop string, or the loop cannot terminate.

    Several WARN rules are heuristics a legitimate workbook can never clear (a
    genuinely header-less matrix; the ±2 tolerance of the rendered-width
    estimate). Gating delivery on them once produced ``结论: 有条件通过`` — a tier
    that reads as helpful and is really an infinite loop, because the model was
    told to iterate until ``结论: 通过``.
    """
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "名单"
    sheet.append(["姓名", "部门", "工号"])
    sheet.append(["张三", "研发", "A001"])
    sheet.append(["李四", "待填", "A002"])  # residual placeholder -> WARN, no formulas -> no ERROR
    path = tmp_path / "warn_only.xlsx"
    workbook.save(path)

    out = run_script("bisheng-xlsx", "inspect_workbook.py", path, "--check-only").stdout

    total = re.search(r"合计: (\d+) ERROR / (\d+) WARN", out)
    assert total, out
    assert total.group(1) == "0", out
    assert int(total.group(2)) >= 1, f"sample was meant to trip a WARN:\n{out}"
    assert PASS in out, out
