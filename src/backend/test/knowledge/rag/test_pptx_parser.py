from types import SimpleNamespace

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from bisheng.pptx2md.parser import has_meaningful_shape_text, is_text_block
from bisheng.pptx2md.types import ConversionConfig


def make_config():
    return ConversionConfig(
        pptx_path="dummy.pptx",
        output_path="dummy.md",
        image_dir=".",
    )


def make_shape(
    *,
    text: str,
    shape_type,
    has_text_frame: bool = True,
    is_placeholder: bool = False,
    placeholder_type=None,
):
    placeholder_format = SimpleNamespace(type=placeholder_type)
    return SimpleNamespace(
        text=text,
        shape_type=shape_type,
        has_text_frame=has_text_frame,
        is_placeholder=is_placeholder,
        placeholder_format=placeholder_format,
    )


def test_has_meaningful_shape_text_filters_decorative_numeric_badges():
    assert not has_meaningful_shape_text("1")
    assert not has_meaningful_shape_text(" | ")
    assert has_meaningful_shape_text("P08 数据治理")
    assert has_meaningful_shape_text("顶层设计共识难")


def test_is_text_block_keeps_short_auto_shape_labels():
    config = make_config()
    shape = make_shape(text="顶层设计共识难", shape_type=MSO_SHAPE_TYPE.AUTO_SHAPE)

    assert is_text_block(config, shape)


def test_is_text_block_keeps_short_freeform_labels():
    config = make_config()
    shape = make_shape(text="推动智改数转\n建设“智慧中粮”", shape_type=MSO_SHAPE_TYPE.FREEFORM)

    assert is_text_block(config, shape)


def test_is_text_block_skips_pure_numeric_auto_shape_badges():
    config = make_config()
    shape = make_shape(text="2", shape_type=MSO_SHAPE_TYPE.AUTO_SHAPE)

    assert not is_text_block(config, shape)


def test_is_text_block_keeps_placeholder_body_text():
    config = make_config()
    shape = make_shape(
        text="通过对现状访谈调研梳理出的1836个问题进行分析",
        shape_type=MSO_SHAPE_TYPE.PLACEHOLDER,
        is_placeholder=True,
        placeholder_type=PP_PLACEHOLDER.BODY,
    )

    assert is_text_block(config, shape)
