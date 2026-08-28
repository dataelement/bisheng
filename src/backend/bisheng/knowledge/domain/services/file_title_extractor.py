"""Extract document titles from various file formats without using LLMs."""
# ruff: noqa: RUF001

import base64
import os
import re
from io import BytesIO

from loguru import logger


class FileTitleExtractError(Exception):
    """Raised when title extraction fails for a specific file."""

    pass


# -----------------------------------------------------------------------------
# Generic header / weak-candidate filters
# -----------------------------------------------------------------------------

_GENERIC_HEADERS = frozenset(
    {
        "中华人民共和国国家标准",
        "中华人民共和国公共安全行业标准",
        "国家标准",
        "行业标准",
        "企业标准",
        "团体标准",
        "地方标准",
        "前言",
        "目次",
        "目录",
        "范围",
        "封面",
        "发布",
        "实施",
        "ICS",
        "CCS",
        "UDC",
        "备案号",
    }
)

_INSTITUTION_SUFFIXES = (
    "公司",
    "集团",
    "总公司",
    "分公司",
    "子公司",
    "局",
    "厅",
    "部",
    "委",
    "院",
    "所",
    "大学",
    "学院",
    "委员会",
    "协会",
    "学会",
    "联合会",
    "研究中心",
    "研究室",
    "办公室",
    "中心",
)

_DEFAULT_METADATA_PATTERNS = (
    r"^microsoft\s+(word|excel|powerpoint|office)",
    r"^wps\s*(文字|表格|演示|office)?",
    r"^新建\s*文档?$",
    r"^未命名\s*文档?$",
    r"^新建\s*microsoft\s+(word|excel|powerpoint)",
    r"^文档\d*$",
    r"^演示文稿\d*$",
    r"^工作簿\d*$",
    r"^document\s*\d*$",
    r"^presentation\s*\d*$",
    r"^workbook\s*\d*$",
    r"^book\s*\d*$",
    r"^template\s*\d*$",
    r"^draft\s*\d*$",
    r"^scan\s*\d*$",
    r"^img_\d+",
    r"^image_\d+",
    r"^untitled\s*\d*$",
)


_PATH_LIKE_RE = re.compile(
    r"[\\/:]|^[A-Za-z]:\\|\.exe\b|\.tmp\b|\.temp\b|file://|ftp://|http://|https://|www\.|[\w.-]+@[\w.-]+\.\w+"
)


def _clean_text(text: str) -> str:
    """Collapse whitespace and strip a candidate string."""
    return re.sub(r"\s+", " ", text).strip()


def _halfwidth(text: str) -> str:
    """Convert common full-width alphanumerics to half-width for comparison."""
    table = str.maketrans(
        "０１２３４５６７８９ａｂｃｄｅｆｇｈｉｊｋｌｍｎｏｐｑｒｓｔｕｖｗｘｙｚＡＢＣＤＥＦＧＨＩＪＫＬＭＮＯＰＱＲＳＴＵＶＷＸＹＺ",
        "0123456789abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ",
    )
    return text.translate(table)


def _normalize_for_compare(text: str) -> str:
    """Normalize text for equality checks (drop spaces, brackets, separators, case)."""
    if not text:
        return ""
    text = _halfwidth(text).lower()
    text = re.sub(r"[\s\-_《》（）()\[\]【】]", "", text)
    return text


def _has_substantive_overlap(candidate: str, reference: str | None, min_chars: int = 2) -> bool:
    """Return True when ``candidate`` shares a substantive substring with ``reference``."""
    if not reference:
        return False
    ref = _normalize_for_compare(reference)
    cand = _normalize_for_compare(candidate)
    if not ref or not cand:
        return False
    if cand in ref or ref in cand:
        return True
    for i in range(len(cand) - min_chars + 1):
        if cand[i : i + min_chars] in ref:
            return True
    return False


def _is_generic_header(text: str) -> bool:
    """Match common Chinese standard page headers and isolated metadata."""
    t = _clean_text(text)
    if not t:
        return True
    if t in _GENERIC_HEADERS:
        return True
    for h in _GENERIC_HEADERS:
        if t.startswith(h) and len(t) <= len(h) + 3:
            return True
    # Isolated standard number / date / page number
    if re.fullmatch(r"[A-Za-z]{1,6}(?:/[A-Za-z]{1,6})?\s*\d+(?:\.\d+)?\s*[—–\-]\s*\d{2,4}", t):
        return True
    if re.fullmatch(r"\d{4}[年/\-]\d{1,2}[月/\-]?\d{0,2}日?", t):
        return True
    if re.fullmatch(r"第\s*\d+\s*页|\d+\s*/\s*\d+|[-—]\s*\d+\s*[-—]", t):
        return True
    return False


def _is_software_or_path(text: str) -> bool:
    """Detect file paths, URLs, emails, temporary file names and software hints."""
    t = text.lower()
    if _PATH_LIKE_RE.search(t):
        return True
    if any(k in t for k in (".com", ".cn", ".org", ".net", ".gov")):
        return True
    return False


def _is_default_metadata(text: str) -> bool:
    """Detect default editor metadata such as '新建文档' or 'Presentation1'."""
    t = text.lower()
    for pattern in _DEFAULT_METADATA_PATTERNS:
        if re.search(pattern, t, re.UNICODE):
            return True
    return False


def _is_pure_number_or_date(text: str) -> bool:
    """Return True if the text is only a number, version or date."""
    t = _clean_text(text)
    if not t:
        return True
    if re.fullmatch(r"\d+(?:[.,]\d+)*", t):
        return True
    if re.fullmatch(r"\d{4}[年/\-]\d{1,2}[月/\-]?\d{0,2}日?", t):
        return True
    return False


def _is_institution_name(text: str) -> bool:
    """Detect short strings that are only an organization/department name."""
    t = _clean_text(text)
    if len(t) > 40:
        return False
    if t.endswith(_INSTITUTION_SUFFIXES):
        return True
    if re.search(r"(公司|集团|局|厅|部|委|院|所|大学|学院|委员会|协会|学会|办公室|中心)$", t):
        return True
    return False


def _is_weak_attachment_prefix(text: str, file_name: str | None = None) -> bool:
    """Detect leading auxiliary text such as 'Attachment 1 prefix'."""
    t = _clean_text(text)
    if file_name and _has_substantive_overlap(t, os.path.splitext(file_name)[0]):
        return False
    if re.match(r"^附件[一二三四五六七八九十0-9]*[：:.、\-—\s]+", t):
        return True
    return False


def _is_likely_field_header(text: str) -> bool:
    """Detect Excel header rows / table column names."""
    t = _clean_text(text)
    field_words = (
        "序号",
        "名称",
        "项目",
        "内容",
        "备注",
        "合计",
        "总计",
        "数量",
        "单价",
        "金额",
        "日期",
        "时间",
        "编号",
        "编码",
        "类别",
        "类型",
        "状态",
        "操作",
        "单位",
        "部门",
        "人员",
        "姓名",
        "性别",
        "年龄",
        "地址",
        "电话",
        "邮箱",
    )
    if any(w in t for w in field_words) and len(t) < 20:
        return True
    return False


def _is_weak_title(text: str, file_name: str | None = None) -> bool:
    """Return True when ``text`` should not be trusted as a document title."""
    t = _clean_text(text)
    if not t or len(t) < 2 or len(t) > 200:
        return True
    if _is_generic_header(t):
        return True
    if _is_software_or_path(t):
        return True
    if _is_default_metadata(t):
        return True
    if _is_pure_number_or_date(t):
        return True
    if _is_institution_name(t) and not _has_substantive_overlap(t, file_name):
        return True
    if _is_weak_attachment_prefix(t, file_name):
        return True
    return False


def sanitize_file_name(name: str, max_length: int = 200, use_hyphen: bool = False) -> str | None:
    """Clean a candidate file name so it is safe for storage and display.

    Removes leading/trailing whitespace, replaces path separators and other
    illegal characters, collapses multiple spaces, and truncates to the
    configured max length. Returns ``None`` if the sanitized name is empty.
    """
    if not name:
        return None
    name = name.strip()
    name = re.sub(r'[\\/:*?"<>|]', " ", name)
    if use_hyphen:
        name = name.replace("_", "-")
    name = re.sub(r"\s+", " ", name).strip()
    if not name:
        return None
    if len(name) > max_length:
        name = name[:max_length].rstrip()
    return name if name else None


def _read_first_text_block(file_path: str, max_bytes: int = 4096) -> str | None:
    """Read the first non-empty line/paragraph from a plain text file."""
    try:
        with open(file_path, "rb") as f:
            raw = f.read(max_bytes)
        if not raw:
            return None
        for encoding in ("utf-8", "utf-8-sig", "gbk", "gb2312", "latin-1"):
            try:
                text = raw.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            text = raw.decode("utf-8", errors="ignore")
        for line in text.splitlines():
            line = line.strip()
            if line:
                return line
    except Exception as e:
        logger.warning("read_first_text_block failed: {}", e)
    return None


class BaseTitleExtractor:
    """Base class for format-specific title extractors."""

    def extract(self, file_path: str, context: dict | None = None) -> str | None:
        raise NotImplementedError


class TxtTitleExtractor(BaseTitleExtractor):
    """Extract title from plain text files.

    The first non-empty line is only a weak candidate; reject list items,
    dates, URLs, paths and long body sentences.
    """

    _LIST_ITEM_RE = re.compile(r"^[\d一二三四五六七八九十]+[.．、：:]\s+|^[（(]\d+[)）]|^[-•·]\s+", re.UNICODE)
    _DATE_ONLY_RE = re.compile(r"^\d{4}[年/\-]\d{1,2}[月/\-]?\d{0,2}日?$")

    def extract(self, file_path: str, context: dict | None = None) -> str | None:
        file_name = context.get("file_name") if context else None
        first = _read_first_text_block(file_path)
        if not first:
            return None
        first = first.strip()
        if len(first) > 200:
            return None
        if _is_weak_title(first, file_name):
            return None
        if self._LIST_ITEM_RE.match(first):
            return None
        if self._DATE_ONLY_RE.fullmatch(first):
            return None
        if "。" in first and len(first) > 30:
            return None
        return first


class MarkdownTitleExtractor(BaseTitleExtractor):
    """Extract title from Markdown files.

    Strategy:
        1. YAML front matter ``title`` field.
        2. First level-1 heading.
        3. First short text block as fallback.
    Front matter / H1 values are validated against weak-title filters.
    """

    _YAML_TITLE_RE = re.compile(r"^---\s*\n.*?^title:\s*(.+?)\n.*?^---\s*\n", re.M | re.S)
    _H1_RE = re.compile(r"^#\s+(.+)$", re.M)

    def extract(self, file_path: str, context: dict | None = None) -> str | None:
        file_name = context.get("file_name") if context else None
        try:
            with open(file_path, "rb") as f:
                raw = f.read(8192)
            text = raw.decode("utf-8", errors="ignore")
            m = self._YAML_TITLE_RE.search(text)
            if m:
                title = m.group(1).strip().strip('"').strip("'")
                if title and not _is_weak_title(title, file_name):
                    return title
            m = self._H1_RE.search(text)
            if m:
                title = m.group(1).strip()
                if title and not _is_weak_title(title, file_name):
                    return title
            return TxtTitleExtractor().extract(file_path, context)
        except Exception as e:
            logger.warning("markdown title extract failed: {}", e)
        return None


class HtmlTitleExtractor(BaseTitleExtractor):
    """Extract title from HTML files.

    Strategy:
        1. ``<title>`` tag content.
        2. First ``<h1>`` tag content.
        3. Combine both, stripping site/section suffixes from ``<title>``.
    """

    def extract(self, file_path: str, context: dict | None = None) -> str | None:
        file_name = context.get("file_name") if context else None
        try:
            from bs4 import BeautifulSoup

            with open(file_path, "rb") as f:
                raw = f.read(8192)
            text = raw.decode("utf-8", errors="ignore")
            soup = BeautifulSoup(text, "html.parser")

            title_text = soup.title.get_text(strip=True) if soup.title else None
            h1_tag = soup.find("h1")
            h1_text = h1_tag.get_text(strip=True) if h1_tag else None

            candidate = self._pick_html_title(title_text, h1_text)
            if candidate and not _is_weak_title(candidate, file_name):
                return candidate
        except Exception as e:
            logger.warning("html title extract failed: {}", e)
        return None

    @staticmethod
    def _pick_html_title(title_text: str | None, h1_text: str | None) -> str | None:
        if not title_text and not h1_text:
            return None
        if title_text and h1_text:
            if h1_text in title_text:
                # Prefer the segment of <title> that contains the h1 text.
                for sep in (" - ", " | ", " — ", " – ", " _ ", "-", "|"):
                    if sep in title_text:
                        for part in title_text.split(sep):
                            if h1_text in part.strip():
                                return part.strip()
                return h1_text
            if title_text in h1_text:
                return h1_text
            return h1_text
        return title_text or h1_text


class DocxTitleExtractor(BaseTitleExtractor):
    """Extract title from DOCX files.

    Strategy:
        1. Document core properties ``title`` (validated).
        2. Paragraph with style name ``Title`` / ``标题`` (validated).
        3. Merge consecutive, centered/large paragraphs on the first page into a
           single title block, excluding headers, issuing units and dates.
    """

    def extract(self, file_path: str, context: dict | None = None) -> str | None:
        file_name = context.get("file_name") if context else None
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document(file_path)

            core_title = doc.core_properties.title
            if core_title and core_title.strip() and not _is_weak_title(core_title.strip(), file_name):
                return core_title.strip()

            for para in doc.paragraphs:
                style_name = para.style.name if para.style and para.style.name else ""
                if style_name and ("Title" in style_name or "标题" in style_name):
                    text = para.text.strip()
                    if text and not _is_weak_title(text, file_name):
                        return text

            header_texts = set()
            try:
                for section in doc.sections:
                    for header in (section.header, section.first_page_header, section.even_page_header):
                        for p in header.paragraphs:
                            t = p.text.strip()
                            if t:
                                header_texts.add(t)
            except Exception:
                pass

            candidates = []
            max_scan = 40
            for idx, para in enumerate(doc.paragraphs):
                if idx >= max_scan:
                    break
                text = para.text.strip()
                if not text or text in header_texts:
                    continue
                if _is_weak_title(text, file_name):
                    continue
                sizes = [run.font.size.pt for run in para.runs if run.font.size and run.font.size.pt]
                max_size = max(sizes) if sizes else 0
                is_centered = para.alignment == WD_ALIGN_PARAGRAPH.CENTER
                candidates.append({"text": text, "size": max_size, "centered": is_centered, "idx": idx})

            if not candidates:
                return None

            blocks = []
            current = [candidates[0]]
            for c in candidates[1:]:
                last = current[-1]
                if (
                    c["idx"] == last["idx"] + 1
                    and c["centered"]
                    and last["centered"]
                    and abs(c["size"] - last["size"]) <= 4
                ):
                    current.append(c)
                else:
                    blocks.append(current)
                    current = [c]
            blocks.append(current)

            best_block = None
            best_score = -1
            for block in blocks:
                text = " ".join(p["text"] for p in block)
                size = max(p["size"] for p in block)
                centered = all(p["centered"] for p in block)
                score = size * 3 + len(text) * 0.5
                if centered:
                    score += 10
                if score > best_score:
                    best_score = score
                    best_block = text
            return best_block
        except Exception as e:
            logger.warning("docx title extract failed: {}", e)
        return None


class DocTitleExtractor(BaseTitleExtractor):
    """Extract title from legacy DOC files.

    Strategy:
        1. Try python-docx directly (some .doc files are readable).
        2. Convert to DOCX via LibreOffice and reuse DocxTitleExtractor.
    """

    def extract(self, file_path: str, context: dict | None = None) -> str | None:
        try:
            result = DocxTitleExtractor().extract(file_path, context)
            if result:
                return result
        except Exception:
            pass
        try:
            from bisheng.knowledge.rag.pipeline.loader.utils.libreoffice_converter import (
                convert_doc_to_docx,
            )

            docx_path = convert_doc_to_docx(file_path)
            if docx_path and os.path.exists(docx_path):
                return DocxTitleExtractor().extract(docx_path, context)
        except Exception as e:
            logger.warning("doc title extract failed: {}", e)
        return None


class PptxTitleExtractor(BaseTitleExtractor):
    """Extract title from PPTX files.

    Strategy:
        1. Title placeholder on the first slide.
        2. Merge related text boxes on the first slide, excluding logos,
           department names, footers and boilerplate such as '汇报材料'.
    """

    def extract(self, file_path: str, context: dict | None = None) -> str | None:
        file_name = context.get("file_name") if context else None
        try:
            from pptx import Presentation
            from pptx.enum.shapes import PP_PLACEHOLDER

            prs = Presentation(file_path)
            if not prs.slides:
                return None
            slide = prs.slides[0]
            slide_height = prs.slide_height

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape.placeholder_format is not None and shape.placeholder_format.type in (
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.CENTER_TITLE,
                ):
                    text = shape.text_frame.text.strip()
                    if text and not _is_weak_title(text, file_name):
                        return text

            candidates = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text or _is_weak_title(text, file_name):
                    continue
                lower = text.lower()
                if any(k in lower for k in ("汇报材料", "logo", "页脚", "footer", "部门", "department")):
                    continue
                if shape.top > slide_height * 0.85:
                    continue
                sizes = []
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            sizes.append(run.font.size.pt)
                max_size = max(sizes) if sizes else 0
                candidates.append(
                    {
                        "text": text,
                        "size": max_size,
                        "top": shape.top,
                        "left": shape.left,
                        "width": shape.width,
                    }
                )

            if not candidates:
                return None

            candidates.sort(key=lambda x: x["top"])
            tol = int(slide_height * 0.05) or 1
            groups = []
            current = [candidates[0]]
            for c in candidates[1:]:
                if abs(c["top"] - current[0]["top"]) <= tol:
                    current.append(c)
                else:
                    groups.append(current)
                    current = [c]
            groups.append(current)

            best_text = None
            best_score = -1
            for group in groups:
                group.sort(key=lambda x: x["left"])
                merged = " ".join(s["text"] for s in group)
                max_size = max(s["size"] for s in group)
                top = group[0]["top"]
                score = max_size * 3 + len(merged) * 0.5 - (top / slide_height) * 10
                if score > best_score:
                    best_score = score
                    best_text = merged
            return best_text
        except Exception as e:
            logger.warning("pptx title extract failed: {}", e)
        return None


class PptTitleExtractor(BaseTitleExtractor):
    """Extract title from legacy PPT files.

    Strategy:
        1. Try python-pptx directly.
        2. Convert to PPTX via LibreOffice and reuse PptxTitleExtractor.
    """

    def extract(self, file_path: str, context: dict | None = None) -> str | None:
        try:
            result = PptxTitleExtractor().extract(file_path, context)
            if result:
                return result
        except Exception:
            pass
        try:
            from bisheng.knowledge.rag.pipeline.loader.utils.libreoffice_converter import (
                convert_ppt_to_pptx,
            )

            pptx_path = convert_ppt_to_pptx(file_path)
            if pptx_path and os.path.exists(pptx_path):
                return PptxTitleExtractor().extract(pptx_path, context)
        except Exception as e:
            logger.warning("ppt title extract failed: {}", e)
        return None


class ExcelTitleExtractor(BaseTitleExtractor):
    """Extract title from XLS/XLSX files.

    Strategy:
        1. Workbook properties title (validated).
        2. Top merged cell value on the first worksheet (validated, not a header row).
        3. Cell A1 only when it does not look like a column header.
    """

    def extract(self, file_path: str, context: dict | None = None) -> str | None:
        file_name = context.get("file_name") if context else None
        try:
            from openpyxl import load_workbook

            path = file_path
            if file_path.lower().endswith(".xls"):
                from bisheng.knowledge.rag.pipeline.loader.utils.md_from_excel import (
                    xls_to_xlsx,
                )

                converted = xls_to_xlsx(file_path)
                if converted and os.path.exists(converted):
                    path = converted
                else:
                    return None

            wb = load_workbook(path, read_only=True, data_only=True)

            if wb.properties and wb.properties.title:
                title = wb.properties.title.strip()
                if title and not _is_weak_title(title, file_name):
                    return title

            if not wb.sheetnames:
                return None
            ws = wb[wb.sheetnames[0]]

            if ws.merged_cells.ranges:
                for merged_range in ws.merged_cells.ranges:
                    min_col, min_row, _, _ = merged_range.bounds
                    if min_row <= 3:
                        value = ws.cell(row=min_row, column=min_col).value
                        if value:
                            text = str(value).strip()
                            if text and not _is_weak_title(text, file_name) and not _is_likely_field_header(text):
                                return text

            value = ws.cell(row=1, column=1).value
            if value:
                text = str(value).strip()
                if text and not _is_weak_title(text, file_name) and not _is_likely_field_header(text):
                    return text
        except Exception as e:
            logger.warning("excel title extract failed: {}", e)
        return None


class CsvTitleExtractor(BaseTitleExtractor):
    """CSV files usually only contain field names, no document title."""

    def extract(self, _file_path: str, _context: dict | None = None) -> str | None:
        return None


class PdfTitleExtractor(BaseTitleExtractor):
    """Extract title from PDF files.

    Strategy:
        1. Metadata title (validated).
        2. Merge first-page spans into title blocks by position and font size.
        3. OCR the top region and merge adjacent OCR lines.
    """

    def __init__(self, ocr_top_ratio: float = 0.25) -> None:
        self.ocr_top_ratio = ocr_top_ratio

    def extract(self, file_path: str, context: dict | None = None) -> str | None:
        file_name = context.get("file_name") if context else None
        try:
            import fitz  # pymupdf

            doc = fitz.open(file_path)
            metadata = doc.metadata or {}
            title = metadata.get("title")
            if title and title.strip() and not _is_weak_title(title.strip(), file_name):
                return title.strip()

            if not doc:
                return None
            page = doc[0]

            title_block = self._extract_title_block(page, file_name)
            if title_block:
                return title_block

            ocr_title = self._ocr_top_region(page, file_name)
            if ocr_title:
                return ocr_title
        except Exception as e:
            logger.warning("pdf title extract failed: {}", e)
        return None

    def _extract_title_block(self, page, file_name: str | None = None) -> str | None:
        try:
            blocks = page.get_text("dict").get("blocks", [])
            page_height = page.rect.height
            page_width = page.rect.width
            spans = []
            for b in blocks:
                if "lines" not in b:
                    continue
                for line in b["lines"]:
                    for span in line["spans"]:
                        text = span.get("text", "").strip()
                        if not text:
                            continue
                        if _is_weak_title(text, file_name):
                            continue
                        bbox = span.get("bbox")
                        if not bbox:
                            continue
                        x0, y0, x1, y1 = bbox
                        if y0 > page_height * 0.6:
                            continue
                        spans.append(
                            {
                                "x0": x0,
                                "y0": y0,
                                "x1": x1,
                                "y1": y1,
                                "size": span.get("size", 0),
                                "text": text,
                            }
                        )
            if not spans:
                return None

            spans.sort(key=lambda s: s["y0"])
            max_size = max(s["size"] for s in spans)
            y_tol = max(2.0, max_size * 0.2)
            lines = []
            current = [spans[0]]
            cy = spans[0]["y0"]
            for s in spans[1:]:
                if abs(s["y0"] - cy) <= y_tol:
                    current.append(s)
                else:
                    lines.append(sorted(current, key=lambda s: s["x0"]))
                    current = [s]
                    cy = s["y0"]
            lines.append(sorted(current, key=lambda s: s["x0"]))

            line_records = []
            for line in lines:
                text = "".join(s["text"] for s in line).strip()
                if not text:
                    continue
                x0 = min(s["x0"] for s in line)
                x1 = max(s["x1"] for s in line)
                y0 = min(s["y0"] for s in line)
                y1 = max(s["y1"] for s in line)
                size = max(s["size"] for s in line)
                line_records.append({"text": text, "x0": x0, "x1": x1, "y0": y0, "y1": y1, "size": size})

            if not line_records:
                return None

            blocks2 = []
            current = [line_records[0]]
            for lr in line_records[1:]:
                last = current[-1]
                gap = lr["y0"] - last["y1"]
                max_h = max(lr["y1"] - lr["y0"], last["y1"] - last["y0"], 1)
                size_ratio = lr["size"] / last["size"] if last["size"] else 1
                overlap = not (lr["x1"] < last["x0"] - 20 or lr["x0"] > last["x1"] + 20)
                left_close = abs(lr["x0"] - last["x0"]) <= 30
                if gap <= max_h * 2.0 and 0.6 <= size_ratio <= 1.7 and (overlap or left_close):
                    current.append(lr)
                else:
                    blocks2.append(current)
                    current = [lr]
            blocks2.append(current)

            best_block = None
            best_score = -1
            for block in blocks2:
                text = " ".join(line["text"] for line in block)
                size = max(line["size"] for line in block)
                min_y = block[0]["y0"]
                char_count = len(text)
                x0 = min(line["x0"] for line in block)
                x1 = max(line["x1"] for line in block)
                centered = abs((x0 + x1) / 2 - page_width / 2) < page_width * 0.15
                score = (
                    size * 5
                    + (1 - min_y / page_height) * 40
                    + min(char_count, 80)
                    + (15 if 5 <= char_count <= 120 else 0)
                    + (10 if centered else 0)
                )
                if score > best_score:
                    best_score = score
                    best_block = text
            return best_block
        except Exception as e:
            logger.warning("pdf title block extraction failed: {}", e)
        return None

    def _ocr_top_region(self, page, file_name: str | None = None) -> str | None:
        try:
            import fitz  # pymupdf

            from bisheng.core.config.settings import settings

            ocr_conf = settings.knowledge.paddle_ocr
            url = ocr_conf.url.strip() if ocr_conf.url else ""
            if not url:
                return None

            clip = fitz.Rect(
                0,
                0,
                page.rect.width,
                page.rect.height * self.ocr_top_ratio,
            )
            pix = page.get_pixmap(clip=clip, dpi=150)
            img_bytes = pix.tobytes("png")
            b64_data = base64.b64encode(img_bytes).decode("utf-8")

            from bisheng.knowledge.rag.pipeline.loader.paddle_ocr import PaddleOcrLoader

            loader = PaddleOcrLoader(
                url=url,
                auth_token=ocr_conf.auth_token or None,
                headers=ocr_conf.headers or None,
                timeout=ocr_conf.timeout or 60,
            )
            result = loader._call_api_sync(b64_data)
            layout_results = result.get("layoutAnalysisResult", [])
            items = loader._extract_parsing_items(layout_results)
            return _merge_ocr_title_items(items, file_name)
        except Exception as e:
            logger.warning("pdf ocr title extract failed: {}", e)
        return None


def _merge_ocr_title_items(items: list[dict], file_name: str | None = None) -> str | None:
    """Merge adjacent OCR items into title blocks for scanned pages/images."""
    if not items:
        return None

    filtered = []
    for item in items:
        text = item.get("text", "").strip()
        if not text or _is_weak_title(text, file_name):
            continue
        bbox = item.get("bbox", [])
        if len(bbox) < 4:
            continue
        try:
            x1, y1, x2, y2 = map(float, bbox[:4])
        except (ValueError, TypeError):
            continue
        filtered.append(
            {
                "text": text,
                "x1": x1,
                "y1": y1,
                "x2": x2,
                "y2": y2,
                "height": y2 - y1,
                "type": item.get("type", "text"),
            }
        )
    if not filtered:
        return None

    filtered.sort(key=lambda i: i["y1"])
    avg_height = sum(i["height"] for i in filtered) / len(filtered)
    y_tol = max(5.0, avg_height * 0.5)
    lines = []
    current = [filtered[0]]
    cy = filtered[0]["y1"]
    for it in filtered[1:]:
        if abs(it["y1"] - cy) <= y_tol:
            current.append(it)
        else:
            lines.append(sorted(current, key=lambda i: i["x1"]))
            current = [it]
            cy = it["y1"]
    lines.append(sorted(current, key=lambda i: i["x1"]))

    line_records = []
    for line in lines:
        text = "".join(i["text"] for i in line).strip()
        if not text:
            continue
        x1 = min(i["x1"] for i in line)
        x2 = max(i["x2"] for i in line)
        y1 = min(i["y1"] for i in line)
        y2 = max(i["y2"] for i in line)
        height = max(i["height"] for i in line)
        is_title = any(i["type"] == "Title" for i in line)
        line_records.append(
            {"text": text, "x1": x1, "x2": x2, "y1": y1, "y2": y2, "height": height, "is_title": is_title}
        )

    if not line_records:
        return None

    blocks = []
    current = [line_records[0]]
    for lr in line_records[1:]:
        last = current[-1]
        gap = lr["y1"] - last["y2"]
        max_h = max(lr["height"], last["height"], 1)
        x_overlap = not (lr["x2"] < last["x1"] - 20 or lr["x1"] > last["x2"] + 20)
        left_close = abs(lr["x1"] - last["x1"]) <= 30
        if gap <= max_h * 2.0 and (x_overlap or left_close):
            current.append(lr)
        else:
            blocks.append(current)
            current = [lr]
    blocks.append(current)

    max_y = max(lr["y2"] for block in blocks for lr in block)
    best = None
    best_score = -1
    for block in blocks:
        text = " ".join(line["text"] for line in block)
        height = max(line["height"] for line in block)
        y1 = block[0]["y1"]
        is_title = any(line["is_title"] for line in block)
        char_count = len(text)
        score = (
            height * 5
            + (1 - y1 / max(max_y, 1)) * 40
            + min(char_count, 80)
            + (20 if is_title else 0)
            + (15 if 5 <= char_count <= 120 else 0)
        )
        if score > best_score:
            best_score = score
            best = text
    return best or None


class ImageTitleExtractor(BaseTitleExtractor):
    """Extract title from images via OCR.

    Strategy:
        1. Crop the top region of the image.
        2. Call PaddleOCR HTTP API.
        3. Merge adjacent OCR lines into title blocks instead of taking the
           topmost fragment.
    """

    def __init__(self, ocr_top_ratio: float = 0.25) -> None:
        self.ocr_top_ratio = ocr_top_ratio

    def extract(self, file_path: str, context: dict | None = None) -> str | None:
        file_name = context.get("file_name") if context else None
        try:
            from bisheng.core.config.settings import settings

            ocr_conf = settings.knowledge.paddle_ocr
            url = ocr_conf.url.strip() if ocr_conf.url else ""
            if not url:
                return None

            from PIL import Image

            with Image.open(file_path) as img:
                width, height = img.size
                top_height = max(1, int(height * self.ocr_top_ratio))
                crop_box = (0, 0, width, top_height)
                cropped = img.crop(crop_box)
                buffer = BytesIO()
                cropped.save(buffer, format="PNG")
                b64_data = base64.b64encode(buffer.getvalue()).decode("utf-8")

            from bisheng.knowledge.rag.pipeline.loader.paddle_ocr import PaddleOcrLoader

            loader = PaddleOcrLoader(
                url=url,
                auth_token=ocr_conf.auth_token or None,
                headers=ocr_conf.headers or None,
                timeout=ocr_conf.timeout or 60,
            )
            result = loader._call_api_sync(b64_data)
            layout_results = result.get("layoutAnalysisResult", [])
            items = loader._extract_parsing_items(layout_results)
            return _merge_ocr_title_items(items, file_name)
        except Exception as e:
            logger.warning("image title extract failed: {}", e)
        return None


class FileTitleExtractorService:
    """Facade for extracting document titles from supported file formats."""

    _EXTRACTORS: dict[str, BaseTitleExtractor] = {
        "txt": TxtTitleExtractor(),
        "md": MarkdownTitleExtractor(),
        "html": HtmlTitleExtractor(),
        "docx": DocxTitleExtractor(),
        "doc": DocTitleExtractor(),
        "pptx": PptxTitleExtractor(),
        "ppt": PptTitleExtractor(),
        "xlsx": ExcelTitleExtractor(),
        "xls": ExcelTitleExtractor(),
        "csv": CsvTitleExtractor(),
        "pdf": PdfTitleExtractor(),
        "png": ImageTitleExtractor(),
        "jpg": ImageTitleExtractor(),
        "jpeg": ImageTitleExtractor(),
    }

    @classmethod
    def extract_title(cls, file_path: str, original_file_name: str | None = None) -> str | None:
        """Extract the title from *file_path* based on its extension.

        Returns the raw title string or ``None`` if no title could be extracted.
        The caller is responsible for sanitizing the result before using it as a
        file name.
        """
        if not file_path or not os.path.exists(file_path):
            logger.info("title extraction skipped, file missing file_path={}", file_path)
            return None
        ext = os.path.splitext(file_path)[1].lower().lstrip(".")
        extractor = cls._EXTRACTORS.get(ext)
        logger.info(
            "title extraction dispatch file_path={} extension={} extractor={}",
            file_path,
            ext,
            type(extractor).__name__ if extractor else None,
        )
        if extractor is None:
            logger.info("no title extractor for extension: {}", ext)
            return None
        try:
            context = {"file_name": original_file_name} if original_file_name else None
            title = extractor.extract(file_path, context=context)
            logger.info(
                "title extraction done file_path={} extension={} title={}",
                file_path,
                ext,
                title,
            )
            return title
        except Exception as e:
            logger.warning("title extraction failed for {}: {}", file_path, e)
            return None
