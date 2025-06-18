# Copyright 2024 Liu Siyao
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import logging
from operator import attrgetter
from typing import Optional

import numpy as np
import pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Length
from scipy.optimize import curve_fit

from bisheng.pptx2md.types import MultiColumnSlide

logger = logging.getLogger(__name__)


def normal_pdf(x_vector, mu=0, sigma=1):
    return (1 / (sigma * np.sqrt(2 * np.pi))) * np.exp(-((x_vector - mu) / sigma)**2 / 2)


def f(x_vector, theta0, theta1, sigma0, sigma1):
    # sigma = 100
    return 0.5 * ((1 / (sigma0 * np.sqrt(2 * np.pi))) * np.exp(-((x_vector - theta0) / sigma0)**2 / 2) +
                  (1 / (sigma1 * np.sqrt(2 * np.pi))) * np.exp(-((x_vector - theta1) / sigma1)**2 / 2))


def f_gauss1(x_vector, theta0, sigma0):
    salida = normal_pdf(x_vector, theta0, sigma0)
    return (salida)


def f_gauss2(x_vector, theta0, theta1, sigma0, sigma1):
    salida = (normal_pdf(x_vector, theta0, sigma0) + normal_pdf(x_vector, theta1, sigma1)) / 2
    return (salida)


def f_gauss3(x_vector, theta0, theta1, theta2, sigma0, sigma1, sigma2):
    salida = (normal_pdf(x_vector, theta0, sigma0) + normal_pdf(x_vector, theta1, sigma1) +
              normal_pdf(x_vector, theta2, sigma2)) / 3
    return (salida)


def compute_pdf_overlap(pdf_fun1, pdf_fun2):
    fun_array = np.vstack([pdf_fun1, pdf_fun2])
    intersection = np.min(fun_array, axis=0)
    pdf_overlap = np.sum(intersection)
    return (pdf_overlap)


def fit_column_model(x_val, g_val):

    q1 = np.quantile(x_val, 0.25)
    q2 = np.median(x_val)
    q3 = np.quantile(x_val, 0.75)

    # print("Using q1: %d, q2: %d, q3: %d"%(q1, q2, q3))

    try:
        params1, cov1 = curve_fit(f_gauss1, x_val, g_val, [q2, q2 - q1])
    except:
        params1 = [q2, q2 - 1]

    try:
        params2, cov2 = curve_fit(f_gauss2, x_val, g_val, [q1, q3, q1, q1])
    except:
        params2 = [q1, q3, q1, q1]

    try:
        params3, cov3 = curve_fit(f_gauss3, x_val, g_val, [q1, q2, q3, q1, q2 - q1, q1])
    except:
        params3 = [q1, q2, q3, q1, q2 - q1, q1]

    # Extract area under the curve of the intersection
    auc1 = compute_pdf_overlap(f_gauss1(x_val, *params1), g_val)
    auc2 = compute_pdf_overlap(f_gauss2(x_val, *params2), g_val)
    auc3 = compute_pdf_overlap(f_gauss3(x_val, *params3), g_val)

    print("Using auc1: %.2f, auc2: %.2f, auc3: %.2f" % (auc1, auc2, auc3))

    if auc1 > 0.86:
        print("Selected 1")
        return (params1)
    elif auc2 > 0.86:
        print("Selected 2")
        return (params2)
    elif auc3 > 0.86:
        print("Selected 3")
        return (params3)
    else:
        idx = np.argmax([auc1, auc2, auc3])
        all_params = [params1, params2, params3]
        print("Selected %d" % (idx + 1))
        return (all_params[idx])


def ungroup_shapes(shapes):
    res = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                res.extend(ungroup_shapes(shape.shapes))
            else:
                res.append(shape)
        except Exception as e:
            print(f'failed to load shape {shape}, skipped. error: {e}')
    return res


def is_two_column_text(slide):

    if slide.slide_layout.name != "TITLE":
        all_mu = list()
        all_sigma = list()
        for shape in sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left')):
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    if shape.has_text_frame:
                        print('SLIDE TITLE: %s' % shape.text_frame.text)

                    continue

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.has_text_frame:
                centroid_x = shape.left + shape.width / 2
                all_mu.append(Length(centroid_x).mm)
                all_sigma.append(Length(shape.width / 4).mm)  # Gaussiana - 2sigma

        return (all_mu, all_sigma)
    else:
        return False


def assign_shapes(slide, params, ncols=2, slide_width_mm=1000):

    shapes_dict = {"shapes_pre": list(), "shapes_l": list(), "shapes_c": list(), "shapes_r": list()}

    shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left'))

    print("Ncols is %d" % ncols)

    if ncols == 1:
        shapes_dict["shapes_pre"] = sorted(shapes,
                                           key=lambda x: x.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER,
                                           reverse=True)
        return (shapes_dict)
    elif ncols == 2:
        param_means = params[0:2]
        param_sds = params[2:]
    elif ncols == 3:
        param_means = params[0:3]
        param_sds = params[3:]
    else:
        raise (ValueError, "Error in the number of columns")

    x_vector = np.arange(1, slide_width_mm)

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                if shape.has_text_frame:
                    print('SLIDE TITLE: %s' % shape.text_frame.text)
                    shapes_dict["shapes_pre"].insert(0, shape)
                else:
                    shapes_dict["shapes_pre"].append(shape)
                continue

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.has_text_frame:
            centroid_x = shape.left + shape.width / 2
            curr_mu = Length(centroid_x).mm
            curr_sigma = Length(shape.width / 4).mm  # Gaussian - 2sigma

            area_u_c = np.zeros(ncols)

            for idx, param_mu in enumerate(param_means):
                area_u_c[idx] = compute_pdf_overlap(normal_pdf(x_vector, mu=param_mu, sigma=param_sds[idx]),
                                                    normal_pdf(x_vector, curr_mu, curr_sigma))

            max_score_column = np.argmax(area_u_c)

            if max_score_column == 0:
                shapes_dict["shapes_l"].append(shape)
            elif max_score_column == 1:
                if ncols == 2:
                    shapes_dict["shapes_r"].append(shape)
                elif ncols == 3:
                    shapes_dict["shapes_c"].append(shape)
                else:
                    raise (ValueError, "Not allowed number of columns")
            elif max_score_column == 2:
                shapes_dict["shapes_r"].append(shape)
            else:
                raise (ValueError, "Max number of columns does not correspond to the number of columns")

    return (shapes_dict)


def get_multi_column_slide_if_present(prs: Presentation, raw_slide, process_shapes) -> Optional[MultiColumnSlide]:
    pdf_modelo = is_two_column_text(raw_slide)

    if not pdf_modelo:
        return None

    slide_width_mm = pptx.util.Length(prs.slide_width)
    t_vector = np.arange(1, slide_width_mm)

    # Model to infer number of columns
    salida = map(lambda mu, sigma: normal_pdf(t_vector, mu, sigma), pdf_modelo[0], pdf_modelo[1])
    sum_of_gaussian = np.mean(list(salida), axis=0)
    parameters = fit_column_model(t_vector, sum_of_gaussian)

    num_cols = int(len(parameters) / 2)

    if num_cols == 1:
        return None

    slide = MultiColumnSlide(preface=[], columns=[], notes=[])

    dict_shapes = assign_shapes(raw_slide, parameters, num_cols, slide_width_mm=slide_width_mm)

    slide.preface = process_shapes(dict_shapes["shapes_pre"])

    if num_cols == 2:
        slide.columns = [process_shapes(dict_shapes["shapes_l"]), process_shapes(dict_shapes["shapes_r"])]

    elif num_cols == 3:
        slide.columns = [
            process_shapes(dict_shapes["shapes_l"]),
            process_shapes(dict_shapes["shapes_c"]),
            process_shapes(dict_shapes["shapes_r"])
        ]

    return slide
