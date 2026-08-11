#!/usr/bin/env python
"""Describe a .pptx/.potx template so its layouts can be reused instead of guessed at.

python-pptx cannot duplicate an existing slide, but it *can* add a slide from any
layout the template defines — inheriting that template's master, theme colours,
fonts and placeholder geometry. To do that you need to know which layout index
does what, and which placeholder idx to fill. That is what this prints.

Usage:
    python skills/bisheng-pptx/scripts/probe_template.py uploads/模板.pptx
    python skills/bisheng-pptx/scripts/probe_template.py uploads/模板.pptx --layout 3
"""

from __future__ import annotations

import argparse
import os
import sys
import traceback

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

EMU_PER_INCH = 914400
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

COLOR_LABELS = {
    "dk1": "深色1/正文色",
    "lt1": "浅色1/背景色",
    "dk2": "深色2",
    "lt2": "浅色2",
    "accent1": "强调色1",
    "accent2": "强调色2",
    "accent3": "强调色3",
    "accent4": "强调色4",
    "accent5": "强调色5",
    "accent6": "强调色6",
    "hlink": "超链接",
    "folHlink": "已访问链接",
}


def _inch(value) -> float:
    return 0.0 if value is None else float(value) / EMU_PER_INCH


def _ph_label(shape) -> str:
    fmt = shape.placeholder_format
    try:
        type_name = str(fmt.type).split(" ")[0]
    except (ValueError, AttributeError):
        type_name = "?"
    return f"idx={fmt.idx} type={type_name}"


def _theme_root(prs):
    """The first master's theme part as an lxml element, or None."""
    try:
        part = prs.slide_masters[0].part.part_related_by(RT.THEME)
    except (KeyError, IndexError):
        return None
    return etree.fromstring(part.blob)


def _color_of(element) -> str:
    if element is None:
        return "-"
    srgb = element.find(f"{A_NS}srgbClr")
    if srgb is not None:
        return srgb.get("val", "-").upper()
    sys_clr = element.find(f"{A_NS}sysClr")
    if sys_clr is not None:
        return (sys_clr.get("lastClr") or sys_clr.get("val") or "-").upper()
    return "-"


def dump_theme(prs) -> None:
    root = _theme_root(prs)
    print("\n=== 主题（新增页会自动继承，不要硬编码覆盖）===")
    if root is None:
        print("  读不到 theme part，按默认 Office 主题处理")
        return

    scheme = root.find(f".//{A_NS}clrScheme")
    if scheme is not None:
        print("  配色:")
        for child in scheme:
            key = etree.QName(child).localname
            print(f"    {COLOR_LABELS.get(key, key):<12} {key:<10} #{_color_of(child)}")

    fonts = root.find(f".//{A_NS}fontScheme")
    if fonts is not None:
        print("  字体:")
        for tag, label in (("majorFont", "标题"), ("minorFont", "正文")):
            node = fonts.find(f"{A_NS}{tag}")
            if node is None:
                continue
            latin = node.find(f"{A_NS}latin")
            ea = node.find(f"{A_NS}ea")
            latin_name = latin.get("typeface", "-") if latin is not None else "-"
            ea_name = ea.get("typeface", "") if ea is not None else ""
            suffix = f" / 东亚: {ea_name}" if ea_name else " / 东亚: 未指定（会回退到正文字体）"
            print(f"    {label}: {latin_name}{suffix}")


def dump_layouts(prs, only: int | None) -> None:
    master = prs.slide_masters[0]
    print(f"\n=== 版式清单（共 {len(master.slide_layouts)} 个）===")
    print("用法: slide = prs.slides.add_slide(prs.slide_layouts[索引])，再按 idx 填占位符")
    for i, layout in enumerate(master.slide_layouts):
        if only is not None and i != only:
            continue
        placeholders = list(layout.placeholders)
        print(f"\n  [{i}] {layout.name}  占位符 {len(placeholders)} 个")
        for ph in placeholders:
            geo = f"位置 {_inch(ph.left):.2f},{_inch(ph.top):.2f}  尺寸 {_inch(ph.width):.2f}×{_inch(ph.height):.2f} in"
            sample = ph.text_frame.text.strip().replace("\n", " ")[:24] if ph.has_text_frame else ""
            sample = f'  示例文字="{sample}"' if sample else ""
            print(f"      {_ph_label(ph):<28} {geo}{sample}")


def dump_slides(prs) -> None:
    slides = list(prs.slides)
    print(f"\n=== 模板自带的页（共 {len(slides)} 页）===")
    if not slides:
        print("  无（纯版式模板，直接 add_slide 即可）")
        return
    print("  这些页会原样留在你的成品里 —— 用不到的必须删掉（见 SKILL.md「删除模板自带页」）")
    for i, slide in enumerate(slides, start=1):
        try:
            layout_name = slide.slide_layout.name
        except (AttributeError, KeyError):
            layout_name = "?"
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().splitlines()[0][:30]
                break
        print(f"  第{i}页  版式={layout_name:<24} 首行文字={title!r}")


def main() -> int:
    parser = argparse.ArgumentParser(description="Probe a pptx/potx template")
    parser.add_argument("pptx", help="模板路径，相对工作区根，例如 uploads/模板.pptx")
    parser.add_argument("--layout", type=int, default=None, help="只看某一个版式的占位符")
    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        print(f"[FATAL] 文件不存在: {args.pptx}（用相对工作区根的路径，上传件通常在 uploads/ 下）")
        return 0

    # Always exit 0: a non-zero exit makes the code interpreter return stderr and
    # drop stdout, which would hide both this report and the reason it stopped.
    try:
        describe(args.pptx, args.layout)
    except Exception:
        print("[FATAL] 解析模板时出错：")
        traceback.print_exc(file=sys.stdout)
        print("常见原因：文件损坏、是加密文档、或其实是老式 .ppt/.pot（python-pptx 只认 .pptx/.potx）。")
        print("→ 请用户另存为 .pptx 后重传；拿不到就放弃套模板，改按 SKILL.md §3 从零创建并告知用户。")
    return 0


def describe(path: str, only_layout: int | None) -> None:
    prs = Presentation(path)
    w_in, h_in = _inch(prs.slide_width), _inch(prs.slide_height)
    ratio = "16:9" if abs(w_in / h_in - 16 / 9) < 0.02 else ("4:3" if abs(w_in / h_in - 4 / 3) < 0.02 else "自定义")
    print(f"=== 画布 ===\n  {w_in:.2f} × {h_in:.2f} in  ({ratio})  —— 不要改画布尺寸，沿用模板的")

    dump_theme(prs)
    dump_layouts(prs, only_layout)
    dump_slides(prs)

    print("\n=== 下一步 ===")
    print("  1) 用 Presentation('该模板路径') 打开它当基底，不要 Presentation() 空开")
    print("  2) 按上面的索引 add_slide，按 idx 定位占位符，赋值 run.text 而不是 text_frame.text")
    print("  3) 删掉模板自带页里用不到的部分，最后另存到 output/")


if __name__ == "__main__":
    sys.exit(main())
