import pytest

from bisheng.knowledge.domain.services.file_title_extractor import (
    FileTitleExtractorService,
    sanitize_file_name,
)


@pytest.fixture
def tmp_file(tmp_path):
    def _make(name: str, content: str | bytes):
        path = tmp_path / name
        if isinstance(content, str):
            path.write_text(content, encoding="utf-8")
        else:
            path.write_bytes(content)
        return str(path)

    return _make


class TestSanitizeFileName:
    def test_removes_illegal_chars(self):
        assert sanitize_file_name('a/b\\c:d*e?f"g<h>i|j') == "a b c d e f g h i j"

    def test_collapses_whitespace(self):
        assert sanitize_file_name("  hello   world  ") == "hello world"

    def test_truncates_long_names(self):
        long_name = "a" * 300
        assert len(sanitize_file_name(long_name)) == 200

    def test_empty_returns_none(self):
        assert sanitize_file_name("") is None
        assert sanitize_file_name("   /   ") is None


class TestTxtExtractor:
    def test_first_line_as_title(self, tmp_file):
        path = tmp_file("doc.txt", "My Document Title\nBody text here.")
        assert FileTitleExtractorService.extract_title(path) == "My Document Title"

    def test_long_first_line_ignored(self, tmp_file):
        path = tmp_file("doc.txt", "a" * 250)
        assert FileTitleExtractorService.extract_title(path) is None


class TestMarkdownExtractor:
    def test_yaml_title(self, tmp_file):
        content = "---\ntitle: YAML Title\n---\n# Heading\n"
        path = tmp_file("doc.md", content)
        assert FileTitleExtractorService.extract_title(path) == "YAML Title"

    def test_first_h1(self, tmp_file):
        path = tmp_file("doc.md", "# Markdown Heading\nSome text.")
        assert FileTitleExtractorService.extract_title(path) == "Markdown Heading"


class TestHtmlExtractor:
    def test_title_tag(self, tmp_file):
        content = "<html><head><title>HTML Title</title></head><body><h1>H1</h1></body></html>"
        path = tmp_file("doc.html", content)
        assert FileTitleExtractorService.extract_title(path) == "H1"

    def test_h1_fallback(self, tmp_file):
        content = "<html><body><h1>Only H1</h1></body></html>"
        path = tmp_file("doc.html", content)
        assert FileTitleExtractorService.extract_title(path) == "Only H1"


class TestDocxExtractor:
    @pytest.mark.skipif(
        pytest.importorskip("docx", reason="python-docx not installed") and False,
        reason="python-docx not installed",
    )
    def test_core_properties_title(self, tmp_file):
        from docx import Document

        doc = Document()
        doc.core_properties.title = "Docx Core Title"
        path = tmp_file("doc.docx", "")
        doc.save(path)
        assert FileTitleExtractorService.extract_title(path) == "Docx Core Title"

    @pytest.mark.skipif(
        pytest.importorskip("docx", reason="python-docx not installed") and False,
        reason="python-docx not installed",
    )
    def test_title_style(self, tmp_file):
        from docx import Document

        doc = Document()
        para = doc.add_paragraph("Styled Title")
        para.style = "Title"
        path = tmp_file("doc.docx", "")
        doc.save(path)
        assert FileTitleExtractorService.extract_title(path) == "Styled Title"


class TestPptxExtractor:
    @pytest.mark.skipif(
        pytest.importorskip("pptx", reason="python-pptx not installed") and False,
        reason="python-pptx not installed",
    )
    def test_title_placeholder(self, tmp_file):
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Pptx Title"
        path = tmp_file("doc.pptx", "")
        prs.save(path)
        assert FileTitleExtractorService.extract_title(path) == "Pptx Title"


class TestExcelExtractor:
    @pytest.mark.skipif(
        pytest.importorskip("openpyxl", reason="openpyxl not installed") and False,
        reason="openpyxl not installed",
    )
    def test_workbook_properties_title(self, tmp_file):
        from openpyxl import Workbook

        wb = Workbook()
        wb.properties.title = "Workbook Title"
        path = tmp_file("doc.xlsx", "")
        wb.save(path)
        assert FileTitleExtractorService.extract_title(path) == "Workbook Title"

    @pytest.mark.skipif(
        pytest.importorskip("openpyxl", reason="openpyxl not installed") and False,
        reason="openpyxl not installed",
    )
    def test_a1_fallback(self, tmp_file):
        from openpyxl import Workbook

        wb = Workbook()
        wb.active["A1"] = "A1 Title"
        path = tmp_file("doc.xlsx", "")
        wb.save(path)
        assert FileTitleExtractorService.extract_title(path) == "A1 Title"


class TestCsvExtractor:
    def test_returns_none(self, tmp_file):
        path = tmp_file("doc.csv", "name,age\nAlice,30\n")
        assert FileTitleExtractorService.extract_title(path) is None


class TestUnsupportedExtension:
    def test_returns_none(self, tmp_file):
        path = tmp_file("doc.unknown", "some content")
        assert FileTitleExtractorService.extract_title(path) is None
