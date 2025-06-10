# Copyright 2024 Liu Siyao
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

from __future__ import print_function

import logging
import os
from functools import partial
from operator import attrgetter
from typing import List, Union

from PIL import Image
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from rapidfuzz import process as fuze_process
from tqdm import tqdm

from bisheng.pptx2md.multi_column import get_multi_column_slide_if_present
from bisheng.pptx2md.types import (
    ConversionConfig,
    GeneralSlide,
    ImageElement,
    ListItemElement,
    ParagraphElement,
    ParsedPresentation,
    SlideElement,
    TableElement,
    TextRun,
    TextStyle,
    TitleElement,
)

logger = logging.getLogger(__name__)

picture_count = 0


def is_title(shape):
    if shape.is_placeholder and (shape.placeholder_format.type == PP_PLACEHOLDER.TITLE or
                                 shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE or
                                 shape.placeholder_format.type == PP_PLACEHOLDER.VERTICAL_TITLE or
                                 shape.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE):
        return True
    return False


def is_text_block(config: ConversionConfig, shape):
    if shape.has_text_frame:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
            return True
        if len(shape.text) > config.min_block_size:
            return True
    return False


def is_list_block(shape) -> bool:
    levels = []
    for para in shape.text_frame.paragraphs:
        if para.level not in levels:
            levels.append(para.level)
        if para.level != 0 or len(levels) > 1:
            return True
    return False


def is_accent(font):
    if font.underline or font.italic or (
            font.color.type == MSO_COLOR_TYPE.SCHEME and
        (font.color.theme_color == MSO_THEME_COLOR.ACCENT_1 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_2 or
         font.color.theme_color == MSO_THEME_COLOR.ACCENT_3 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_4 or
         font.color.theme_color == MSO_THEME_COLOR.ACCENT_5 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_6)):
        return True
    return False


def is_strong(font):
    if font.bold or (font.color.type == MSO_COLOR_TYPE.SCHEME and (font.color.theme_color == MSO_THEME_COLOR.DARK_1 or
                                                                   font.color.theme_color == MSO_THEME_COLOR.DARK_2)):
        return True
    return False


def get_text_runs(para) -> List[TextRun]:
    runs = []
    for run in para.runs:
        result = TextRun(text=run.text, style=TextStyle())
        if result.text == '':
            continue
        try:
            if run.hyperlink.address:
                result.style.hyperlink = run.hyperlink.address
        except:
            result.style.hyperlink = 'error:ppt-link-parsing-issue'
        if is_accent(run.font):
            result.style.is_accent = True
        if is_strong(run.font):
            result.style.is_strong = True
        if run.font.color.type == MSO_COLOR_TYPE.RGB:
            result.style.color_rgb = run.font.color.rgb
        runs.append(result)
    return runs


def process_title(config: ConversionConfig, shape, slide_idx) -> TitleElement:
    text = shape.text_frame.text.strip()
    if config.custom_titles:
        res = fuze_process.extractOne(text, config.custom_titles.keys(), score_cutoff=92)
        if not res:
            return TitleElement(content=text.strip(), level=max(config.custom_titles.values()) + 1)
        else:
            logger.info(f'Title in slide {slide_idx} "{text}" is converted to "{res[0]}" as specified in title file.')
            return TitleElement(content=res[0].strip(), level=config.custom_titles[res[0]])
    else:
        return TitleElement(content=text.strip(), level=1)


def process_text_blocks(config: ConversionConfig, shape, slide_idx) -> List[Union[ListItemElement, ParagraphElement]]:
    results = []
    if is_list_block(shape):
        for para in shape.text_frame.paragraphs:
            if para.text.strip() == '':
                continue
            text = get_text_runs(para)
            results.append(ListItemElement(content=text, level=para.level))
    else:
        # paragraph block
        for para in shape.text_frame.paragraphs:
            if para.text.strip() == '':
                continue
            text = get_text_runs(para)
            results.append(ParagraphElement(content=text))
    return results


def process_picture(config: ConversionConfig, shape, slide_idx) -> Union[ImageElement, None]:
    if config.disable_image:
        return None

    global picture_count

    file_prefix = ''.join(os.path.basename(config.pptx_path).split('.')[:-1])
    pic_name = file_prefix + f'_{picture_count}'
    pic_ext = shape.image.ext
    if not os.path.exists(config.image_dir):
        os.makedirs(config.image_dir)

    output_path = config.image_dir / f'{pic_name}.{pic_ext}'
    common_path = os.path.commonpath([config.output_path, config.image_dir])
    img_outputter_path = os.path.relpath(output_path, common_path)
    with open(output_path, 'wb') as f:
        f.write(shape.image.blob)
        picture_count += 1

    # normal images
    if pic_ext != 'wmf':
        return ImageElement(path=img_outputter_path, width=config.image_width)

    # wmf images, try to convert, if failed, output as original
    try:
        try:
            Image.open(output_path).save(os.path.splitext(output_path)[0] + '.png')
            return ImageElement(path=os.path.splitext(img_outputter_path)[0] + '.png', width=config.image_width)
        except Exception:  # Image failed, try another
            from wand.image import Image
            with Image(filename=output_path) as img:
                img.format = 'png'
                img.save(filename=os.path.splitext(output_path)[0] + '.png')
            logger.info(f'Image {output_path} in slide {slide_idx} converted to png.')
            return ImageElement(path=os.path.splitext(img_outputter_path)[0] + '.png', width=config.image_width)
    except Exception:
        logger.warning(f'Cannot convert wmf image {output_path} in slide {slide_idx} to png, skipped.')
        return None


def process_table(config: ConversionConfig, shape, slide_idx) -> Union[TableElement, None]:
    table = [[sum([get_text_runs(p)
                   for p in cell.text_frame.paragraphs], [])
              for cell in row.cells]
             for row in shape.table.rows]
    if len(table) > 0:
        return TableElement(content=table)
    return None


def ungroup_shapes(shapes) -> List[SlideElement]:
    res = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                res.extend(ungroup_shapes(shape.shapes))
            else:
                res.append(shape)
        except Exception as e:
            logger.warning(f'failed to load shape {shape}, skipped. error: {e}')
    return res


def process_shapes(config: ConversionConfig, current_shapes, slide_id: int) -> List[SlideElement]:
    results = []
    for shape in current_shapes:
        if is_title(shape):
            results.append(process_title(config, shape, slide_id))
        elif is_text_block(config, shape):
            results.extend(process_text_blocks(config, shape, slide_id))
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                pic = process_picture(config, shape, slide_id)
                if pic:
                    results.append(pic)
            except AttributeError as e:
                logger.warning(f'Failed to process picture, skipped: {e}')
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = process_table(config, shape, slide_id)
            if table:
                results.append(table)
        else:
            try:
                ph = shape.placeholder_format
                if ph.type == PP_PLACEHOLDER.OBJECT and hasattr(shape, "image") and getattr(shape, "image"):
                    pic = process_picture(config, shape, slide_id)
                    if pic:
                        results.append(pic)
            except:
                pass

    return results


def parse(config: ConversionConfig, prs: Presentation) -> ParsedPresentation:
    result = ParsedPresentation(slides=[])

    for idx, slide in enumerate(tqdm(prs.slides, desc='Converting slides')):
        if config.page is not None and idx + 1 != config.page:
            continue
        shapes = []
        try:
            shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left'))
        except:
            logger.warning('Bad shapes encountered in this slide. Please check or remove them and try again.')
            logger.warning('shapes:')
            try:
                for sp in slide.shapes:
                    logger.warning(sp.shape_type)
                    logger.warning(sp.top, sp.left, sp.width, sp.height)
            except:
                logger.warning('failed to print all bad shapes.')

        if not config.try_multi_column:
            result_slide = GeneralSlide(elements=process_shapes(config, shapes, idx + 1))
        else:
            multi_column_slide = get_multi_column_slide_if_present(
                prs, slide, partial(process_shapes, config=config, slide_id=idx + 1))
            if multi_column_slide:
                result_slide = multi_column_slide
            else:
                result_slide = GeneralSlide(elements=process_shapes(config, shapes, idx + 1))

        if not config.disable_notes and slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text
            if text:
                result_slide.notes.append(text)

        result.slides.append(result_slide)

    return result
