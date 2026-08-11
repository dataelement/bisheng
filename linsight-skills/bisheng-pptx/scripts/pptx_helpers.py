#!/usr/bin/env python
"""The handful of python-pptx operations that need raw OOXML, done correctly.

Import from a build script running in the code interpreter:

    import sys
    sys.path.insert(0, "skills/bisheng-pptx/scripts")
    from pptx_helpers import set_font, add_bullet, delete_slide, fill_text

Everything here is standard python-pptx plus lxml — no extra dependencies.
"""

from __future__ import annotations

import copy

from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

__all__ = [
    "add_bullet",
    "delete_slide",
    "fill_text",
    "iter_runs",
    "no_bullet",
    "set_font",
    "set_font_all",
]


def set_font(run, name: str, size_pt: float | None = None, bold: bool | None = None, color=None) -> None:
    """Set a run's font so that **Chinese text actually uses it**.

    ``run.font.name`` only writes ``<a:latin>``, which governs Latin characters.
    CJK glyphs follow ``<a:ea>``; without it PowerPoint falls back to the theme's
    East-Asian font and your font choice silently does nothing to the Chinese.
    ``<a:ea>`` must sit immediately after ``<a:latin>`` — appending it anywhere
    else produces a file PowerPoint refuses to open.
    """
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color

    rPr = run._r.get_or_add_rPr()
    latin = rPr.get_or_add_latin()
    latin.set("typeface", name)
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = latin.makeelement(qn("a:ea"), {})
        latin.addnext(ea)
    ea.set("typeface", name)


def iter_runs(shape):
    """Every run in a shape's text frame."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        yield from para.runs


def set_font_all(shape, name: str, size_pt: float | None = None) -> None:
    """Apply set_font to every run of a shape (handy after filling a template)."""
    for run in iter_runs(shape):
        set_font(run, name, size_pt=size_pt)


def add_bullet(paragraph, char: str = "•", font: str = "Arial", indent_in: float = 0.25) -> None:
    """Give a paragraph a real bullet, with a hanging indent.

    Never type "•" into the text itself — layouts that already supply a bullet
    would then render two. ``<a:buChar>`` must be inserted before ``<a:defRPr>``
    to keep the element order valid.

    ``indent_in`` sets the gap between the bullet and the text (marL with a
    negative first-line indent). Without it the glyph sits flush against the
    first character and wrapped lines do not align under the text.
    """
    pPr = paragraph._p.get_or_add_pPr()
    if indent_in:
        offset = int(Inches(indent_in))
        pPr.set("marL", str(offset))
        pPr.set("indent", str(-offset))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        existing = pPr.find(qn(tag))
        if existing is not None:
            pPr.remove(existing)

    bu_font = pPr.makeelement(qn("a:buFont"), {"typeface": font})
    bu_char = pPr.makeelement(qn("a:buChar"), {"char": char})
    anchor = pPr.find(qn("a:defRPr"))
    if anchor is not None:
        anchor.addprevious(bu_font)
        anchor.addprevious(bu_char)
    else:
        pPr.append(bu_font)
        pPr.append(bu_char)


def no_bullet(paragraph) -> None:
    """Suppress an inherited bullet (e.g. a body placeholder used for prose)."""
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buFont"):
        existing = pPr.find(qn(tag))
        if existing is not None:
            pPr.remove(existing)
    if pPr.find(qn("a:buNone")) is None:
        bu_none = pPr.makeelement(qn("a:buNone"), {})
        anchor = pPr.find(qn("a:defRPr"))
        if anchor is not None:
            anchor.addprevious(bu_none)
        else:
            pPr.append(bu_none)


def delete_slide(prs, index: int) -> None:
    """Remove a slide by position, dropping its relationship too.

    python-pptx has no public delete; leaving the rel behind yields a file that
    PowerPoint reports as damaged. Delete from the highest index first when
    removing several, or the positions shift under you.
    """
    slide_ids = prs.slides._sldIdLst
    entries = list(slide_ids)
    entry = entries[index]
    rel_id = entry.get(qn("r:id"))
    prs.part.drop_rel(rel_id)
    slide_ids.remove(entry)


def fill_text(shape, lines) -> None:
    """Replace a shape's text while keeping the template's formatting.

    ``text_frame.text = "..."`` collapses the whole frame into one unstyled run,
    throwing away the size, colour and bullet the template author set. This
    instead rewrites the first run in place and clones that paragraph for any
    further lines.

    ``lines`` may be a string (split on newlines) or a list of strings.
    """
    if isinstance(lines, str):
        lines = lines.split("\n")
    lines = list(lines) or [""]

    tf = shape.text_frame
    first_p = tf.paragraphs[0]._p
    body = first_p.getparent()

    for para in list(body.findall(qn("a:p")))[1:]:
        body.remove(para)

    def _write(p_el, text: str) -> None:
        runs = p_el.findall(qn("a:r"))
        if not runs:
            run = p_el.makeelement(qn("a:r"), {})
            text_el = p_el.makeelement(qn("a:t"), {})
            run.append(text_el)
            p_el.append(run)
            runs = [run]
        for extra in runs[1:]:
            p_el.remove(extra)
        text_el = runs[0].find(qn("a:t"))
        if text_el is None:
            text_el = runs[0].makeelement(qn("a:t"), {})
            runs[0].append(text_el)
        text_el.text = text
        if text != text.strip():
            text_el.set(qn("xml:space"), "preserve")

    _write(first_p, lines[0])
    for line in lines[1:]:
        clone = copy.deepcopy(first_p)
        _write(clone, line)
        body.append(clone)
